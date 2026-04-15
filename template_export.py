#!/usr/bin/env python3
"""
Template-native AWR PPT exporter.

Uses an Oracle .potx template and only Light/white layouts.
Input can be either:
1) Raw AWR JSON (instances/cohort_rollups/database_statistics/...); or
2) Saved app report JSON containing { raw_data, ui_state }.
"""

from __future__ import annotations

import argparse
import base64
import io
import json
import math
import os
import re
import struct
import tempfile
import zipfile
from collections import defaultdict
from dataclasses import dataclass
from datetime import datetime, timedelta
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor


ORACLE_COLORS = {
    "text": RGBColor(0x2A, 0x2F, 0x2F),
    "muted": RGBColor(0x6B, 0x74, 0x7A),
    "accent1": RGBColor(0x04, 0x53, 0x6F),
    "accent2": RGBColor(0x6C, 0x3F, 0x49),
    "accent3": RGBColor(0xC7, 0x46, 0x34),
    "accent4": RGBColor(0xF0, 0xCC, 0x71),
    "accent5": RGBColor(0x89, 0xB2, 0xB0),
    "accent6": RGBColor(0x86, 0xB5, 0x96),
    "panel": RGBColor(0xF8, 0xFA, 0xFC),
    "line": RGBColor(0xD9, 0xDE, 0xDE),
}
PALETTE = [
    ORACLE_COLORS["accent3"],
    ORACLE_COLORS["accent1"],
    ORACLE_COLORS["accent2"],
    ORACLE_COLORS["accent5"],
    ORACLE_COLORS["accent6"],
    ORACLE_COLORS["accent4"],
]
# Conversion used for Exadata-family ECPU sizing (Exascale and Exadata Dedicated).
ECPU_PER_VCPU_EXADATA = 2.0


def short_label(s: Any, n: int = 34) -> str:
    t = str(s or "")
    return t if len(t) <= n else f"{t[: n - 1]}…"


def to_num(v: Any, default: float = 0.0) -> float:
    try:
        return float(v)
    except Exception:
        return default


def fmt_num(v: Any, dec: int = 1) -> str:
    try:
        x = float(v)
    except Exception:
        return "N/A"
    if dec == 0:
        return f"{int(round(x)):,}"
    return f"{x:,.{dec}f}"


def decode_f64_b64(s: str) -> list[float]:
    raw = base64.b64decode(s + "=" * ((4 - (len(s) % 4)) % 4))
    out = []
    for i in range(0, len(raw), 8):
        out.append(struct.unpack("<d", raw[i : i + 8])[0])
    return out


def _plotly_payload_tail(file_content: str) -> str:
    txt = file_content or ""
    i = txt.rfind("Plotly.newPlot(")
    return txt[i:] if i >= 0 else txt


def _extract_plotly_traces(file_content: str) -> list[dict[str, Any]]:
    src = _plotly_payload_tail(file_content)
    if not src:
        return []
    i0 = src.find("(")
    if i0 < 0:
        return []
    c1 = src.find(",", i0 + 1)  # after div id
    if c1 < 0:
        return []
    a0 = src.find("[", c1 + 1)  # data traces array
    if a0 < 0:
        return []

    level = 0
    in_str = False
    esc = False
    a1 = -1
    for idx, ch in enumerate(src[a0:], a0):
        if in_str:
            if esc:
                esc = False
            elif ch == "\\":
                esc = True
            elif ch == '"':
                in_str = False
            continue
        if ch == '"':
            in_str = True
        elif ch == "[":
            level += 1
        elif ch == "]":
            level -= 1
            if level == 0:
                a1 = idx
                break
    if a1 <= a0:
        return []
    try:
        parsed = json.loads(src[a0 : a1 + 1])
    except Exception:
        return []
    return parsed if isinstance(parsed, list) else []


def parse_time_axis_from_plot_html(file_content: str) -> list[str]:
    traces = _extract_plotly_traces(file_content)
    if not traces:
        return []
    x = traces[0].get("x")
    if isinstance(x, list):
        return [str(v) for v in x]
    return []


def parse_numeric_axis_from_plot_html(file_content: str) -> list[float]:
    traces = _extract_plotly_traces(file_content)
    if not traces:
        return []
    y = traces[0].get("y")
    if isinstance(y, dict):
        dtype = str(y.get("dtype") or "").lower()
        bdata = y.get("bdata")
        if isinstance(bdata, str) and dtype == "f8":
            try:
                return decode_f64_b64(bdata)
            except Exception:
                return []
    if isinstance(y, list):
        try:
            return [float(v) for v in y]
        except Exception:
            return []
    return []


def parse_timeseries_from_plot_html(file_content: str) -> tuple[list[str], list[float]] | None:
    xs = parse_time_axis_from_plot_html(file_content)
    ys = parse_numeric_axis_from_plot_html(file_content)
    if not xs or not ys:
        return None
    n = min(len(xs), len(ys))
    return xs[:n], ys[:n]


def quantile(sorted_vals: list[float], q: float) -> float:
    if not sorted_vals:
        return 0.0
    if len(sorted_vals) == 1:
        return sorted_vals[0]
    pos = (len(sorted_vals) - 1) * q
    base = int(math.floor(pos))
    frac = pos - base
    if base + 1 < len(sorted_vals):
        return sorted_vals[base] + frac * (sorted_vals[base + 1] - sorted_vals[base])
    return sorted_vals[base]


def stats_from_values(vals: list[float]) -> dict[str, float] | None:
    clean = sorted([float(v) for v in vals if isinstance(v, (int, float)) and math.isfinite(float(v))])
    if not clean:
        return None
    return {
        "min": clean[0],
        "p30": quantile(clean, 0.30),
        "p50": quantile(clean, 0.50),
        "p70": quantile(clean, 0.70),
        "p95": quantile(clean, 0.95),
        "max": clean[-1],
    }


def pct(numer: float, denom: float) -> float:
    if denom <= 0:
        return 0.0
    return 100.0 * numer / denom


def parse_ai_response_lines(ui_state: dict[str, Any] | None, max_lines: int = 3) -> list[str]:
    if not isinstance(ui_state, dict):
        return []
    ai = ui_state.get("aiChat")
    if not isinstance(ai, dict):
        return []
    text = str(ai.get("response") or "").strip()
    if not text:
        return []
    lines: list[str] = []
    for raw in text.splitlines():
        t = raw.strip()
        if not t:
            continue
        t = re.sub(r"^\s*[-*•]+\s*", "", t)
        t = re.sub(r"^\s*\d+\.\s*", "", t)
        if t:
            lines.append(t)
        if len(lines) >= max_lines:
            break
    return lines


def get_ai_app_payload(ui_state: dict[str, Any] | None) -> dict[str, Any] | None:
    if not isinstance(ui_state, dict):
        return None
    ai = ui_state.get("aiChat")
    if not isinstance(ai, dict):
        return None
    payload = ai.get("appPayload")
    return payload if isinstance(payload, dict) else None


def get_report_meta(payload: dict[str, Any] | None, ui_state: dict[str, Any] | None) -> dict[str, str]:
    src: dict[str, Any] = {}
    if isinstance(payload, dict) and isinstance(payload.get("report_metadata"), dict):
        src = payload.get("report_metadata") or {}
    elif isinstance(ui_state, dict) and isinstance(ui_state.get("reportMeta"), dict):
        src = ui_state.get("reportMeta") or {}
    return {
        "customerName": str(src.get("customerName") or "").strip(),
        "salesRepName": str(src.get("salesRepName") or "").strip(),
        "architectName": str(src.get("architectName") or "").strip(),
        "engineerName": str(src.get("engineerName") or "").strip(),
    }


def pick_payload_comments(payload: dict[str, Any] | None, scope: str, key: str | None = None, max_lines: int = 4) -> list[str]:
    if not isinstance(payload, dict):
        return []
    if scope == "global":
        arr = payload.get("global_comments")
        if isinstance(arr, list):
            return [str(x).strip() for x in arr if str(x).strip()][:max_lines]
        return []
    if scope == "cohort" and key:
        mp = payload.get("cohort_comments")
        if isinstance(mp, dict) and isinstance(mp.get(key), list):
            return [str(x).strip() for x in mp.get(key) if str(x).strip()][:max_lines]
        return []
    if scope == "instance" and key:
        mp = payload.get("instance_comments")
        if isinstance(mp, dict) and isinstance(mp.get(key), list):
            return [str(x).strip() for x in mp.get(key) if str(x).strip()][:max_lines]
        return []
    return []


def payload_infra_line(payload: dict[str, Any] | None) -> str:
    if not isinstance(payload, dict):
        return ""
    infra = payload.get("infrastructure_recommendation")
    if not isinstance(infra, dict):
        return ""
    tier = str(infra.get("recommended_tier") or "").strip()
    pos = infra.get("scale_position")
    if not tier:
        return ""
    rationale = infra.get("rationale")
    rat0 = ""
    if isinstance(rationale, list) and rationale:
        rat0 = str(rationale[0] or "").strip()
    if rat0:
        return f"Infrastructure recommendation: {tier} (scale {pos}). {rat0}"
    return f"Infrastructure recommendation: {tier} (scale {pos})."


def infra_assessment(
    model: "Model",
    payload: dict[str, Any] | None = None,
) -> dict[str, Any]:
    del payload  # deterministic sizing logic only

    # Sizing limits used by this deterministic model.
    # These thresholds are also surfaced as footnotes in the executive slide.
    base_limits = {"ecpu": 256.0, "memory_gb": 512.0, "storage_tb": 80.0, "iops": 8000.0}
    xs_limits = {"ecpu": 200.0, "storage_tb": 100.0}
    # Exadata memory policy:
    # use one DB server memory instead of aggregated two-server memory
    # and cap cohort usage at 70% of that one-server memory.
    exa_one_db_server_memory_gb = 1390.0
    exa_cohort_memory_limit_gb = exa_one_db_server_memory_gb * 0.70
    exa_env_limits = {"ecpu": 1520.0, "memory_gb": exa_cohort_memory_limit_gb, "storage_tb": 5000.0, "iops": 89_000_000.0}
    exa_server_capacity = {"db_ecpu": 380.0, "db_memory_gb": exa_cohort_memory_limit_gb, "storage_tb": 300.0, "storage_iops": 15_000_000.0}

    def fmt_req_limit(req: float, lim: float, unit: str, d: int = 1) -> str:
        return f"{fmt_num(req, d)}/{fmt_num(lim, d)} {unit}"

    selected = model.selected_rows()
    cohorts = model.cohorts()
    total_instances = len(selected)
    total_cohorts = len(cohorts)

    # Base Database fit (instance-by-instance).
    base_db_servers = 0
    base_failures: list[str] = []
    for r in selected:
        cpu = float(r.init_cpu_count or r.logical_cpu_count or 0.0)
        mem_gb = float(r.mem_gb or (r.sga_gb + r.pga_gb) or 0.0)
        st = model.db_metric_stats(r.cohort, r.db, "Allocated Storage (GB)") or {}
        storage_tb = float(st.get("p95") or st.get("max") or 0.0) / 1024.0
        it = model.db_metric_stats(r.cohort, r.db, "DB IOPS") or {}
        iops = float(it.get("p95") or it.get("max") or 0.0)
        db_name = model.display_db_name(r.db)
        inst_label = model.instance_label(r.instance)

        if cpu > base_limits["ecpu"]:
            base_failures.append(f"{db_name} ({inst_label}) CPU {fmt_req_limit(cpu, base_limits['ecpu'], 'ECPU')}")
        if mem_gb > base_limits["memory_gb"]:
            base_failures.append(f"{db_name} ({inst_label}) Memory {fmt_req_limit(mem_gb, base_limits['memory_gb'], 'GB')}")
        if storage_tb > base_limits["storage_tb"]:
            base_failures.append(f"{db_name} ({inst_label}) Storage {fmt_req_limit(storage_tb, base_limits['storage_tb'], 'TB')}")
        if iops > base_limits["iops"]:
            base_failures.append(f"{db_name} ({inst_label}) IOPS {fmt_req_limit(iops, base_limits['iops'], 'IOPS', 0)}")

        base_db_servers += 2 if model.db_is_rac(r.db) else 1

    base_fit = len(base_failures) == 0
    base_reason = "Fit for all instances." if base_fit else f"Not Fit: {base_failures[0]}"

    # Exascale fit (one deployment per cohort).
    xs_failures: list[str] = []
    xs_cohort_details: list[dict[str, Any]] = []
    cohort_matrix: list[dict[str, Any]] = []
    for c in cohorts:
        rs = model.rows_for_cohort(c)
        sm = model.summary(rs)
        c_vcpu = float(sm.get("vcpu_total") or 0.0)
        c_cpu = c_vcpu * ECPU_PER_VCPU_EXADATA
        c_mem_gb = float(sm.get("memory_gb_total") or 0.0)
        c_storage_tb = float(sm.get("allocated_storage_gb") or 0.0) / 1024.0
        c_iops = max((float(v) for _, v in model.metric_by_db(c, "DB IOPS", scale=1.0)), default=0.0)

        largest = sorted(
            rs,
            key=lambda r: ((r.init_cpu_count or r.logical_cpu_count), r.mem_gb, r.sga_gb + r.pga_gb),
            reverse=True,
        )[0] if rs else None
        li_cpu = float((largest.init_cpu_count or largest.logical_cpu_count) if largest else 0.0)
        li_mem_stats = model.db_metric_stats(c, largest.db, "DB Memory (MB)") if largest else None
        li_iops_stats = model.db_metric_stats(c, largest.db, "DB IOPS") if largest else None
        li_storage_stats = model.db_metric_stats(c, largest.db, "Allocated Storage (GB)") if largest else None
        li_mem_gb = float((li_mem_stats or {}).get("p95") or 0.0) / 1024.0
        if li_mem_gb <= 0 and largest:
            li_mem_gb = float((largest.sga_gb or 0.0) + (largest.pga_gb or 0.0))
        li_iops = float((li_iops_stats or {}).get("p95") or (li_iops_stats or {}).get("max") or 0.0)
        li_storage_tb = float((li_storage_stats or {}).get("p95") or (li_storage_stats or {}).get("max") or 0.0) / 1024.0

        base_blockers = []
        if li_cpu > base_limits["ecpu"]:
            base_blockers.append(f"CPU {fmt_req_limit(li_cpu, base_limits['ecpu'], 'ECPU')}")
        if li_storage_tb > base_limits["storage_tb"]:
            base_blockers.append(f"Storage {fmt_req_limit(li_storage_tb, base_limits['storage_tb'], 'TB')}")
        if li_mem_gb > base_limits["memory_gb"]:
            base_blockers.append(f"Memory {fmt_req_limit(li_mem_gb, base_limits['memory_gb'], 'GB')}")
        if li_iops > base_limits["iops"]:
            base_blockers.append(f"IOPS {fmt_req_limit(li_iops, base_limits['iops'], 'IOPS', 0)}")
        base_fit_c = len(base_blockers) == 0

        blockers: list[str] = []
        if c_cpu > xs_limits["ecpu"]:
            blockers.append(f"CPU {fmt_req_limit(c_cpu, xs_limits['ecpu'], 'ECPU')}")
        if c_storage_tb > xs_limits["storage_tb"]:
            blockers.append(f"Storage {fmt_req_limit(c_storage_tb, xs_limits['storage_tb'], 'TB')}")
        xs_cohort_details.append(
            {
                "cohort": c,
                "fit": len(blockers) == 0,
                "reason": "Fit" if len(blockers) == 0 else f"Not Fit: {blockers[0]}",
            }
        )
        if blockers:
            xs_failures.append(f"{c}: {blockers[0]}")

        d_blockers = []
        if c_cpu > exa_env_limits["ecpu"]:
            d_blockers.append(f"CPU {fmt_req_limit(c_cpu, exa_env_limits['ecpu'], 'ECPU')}")
        if c_storage_tb > exa_env_limits["storage_tb"]:
            d_blockers.append(f"Storage {fmt_req_limit(c_storage_tb, exa_env_limits['storage_tb'], 'TB')}")
        if c_mem_gb > exa_env_limits["memory_gb"]:
            d_blockers.append(f"Memory {fmt_req_limit(c_mem_gb, exa_env_limits['memory_gb'], 'GB')}")
        if c_iops > exa_env_limits["iops"]:
            d_blockers.append(f"IOPS {fmt_req_limit(c_iops, exa_env_limits['iops'], 'IOPS', 0)}")

        cohort_matrix.append(
            {
                "cohort": c,
                "largest_instance": largest.instance if largest else "N/A",
                "largest_db_name": model.display_db_name(largest.db) if largest else "N/A",
                "base_fit": base_fit_c,
                "base_reason": "Fit" if base_fit_c else f"Not Fit: {base_blockers[0]}",
                "xs_fit": len(blockers) == 0,
                "xs_reason": "Fit" if len(blockers) == 0 else f"Not Fit: {blockers[0]}",
                "d_fit": len(d_blockers) == 0,
                "d_reason": "Fit" if len(d_blockers) == 0 else f"Not Fit: {d_blockers[0]}",
                "required_ecpu": c_cpu,
                "required_vcpu": c_vcpu,
                "required_storage_tb": c_storage_tb,
            }
        )

    xs_fit = len(xs_failures) == 0
    xs_reason = "Fit for all cohorts." if xs_fit else f"Not Fit: {xs_failures[0]}"

    # Exadata Dedicated inventory for all cohorts consolidated.
    global_s = model.summary()
    req_vcpu = float(global_s.get("vcpu_total") or 0.0)
    req_cpu = req_vcpu * ECPU_PER_VCPU_EXADATA
    req_mem_gb = float(global_s.get("memory_gb_total") or 0.0)
    req_storage_tb = float(global_s.get("allocated_storage_gb") or 0.0) / 1024.0
    req_iops = max((float(v) for _, v in model.iops_by_cohort()), default=0.0)

    exa_envs = max(
        1,
        int(
            math.ceil(
                max(
                    req_cpu / exa_env_limits["ecpu"] if exa_env_limits["ecpu"] > 0 else 0.0,
                    req_mem_gb / exa_env_limits["memory_gb"] if exa_env_limits["memory_gb"] > 0 else 0.0,
                    req_storage_tb / exa_env_limits["storage_tb"] if exa_env_limits["storage_tb"] > 0 else 0.0,
                    req_iops / exa_env_limits["iops"] if exa_env_limits["iops"] > 0 else 0.0,
                )
            )
        ),
    )
    per_env_cpu = req_cpu / exa_envs
    per_env_mem = req_mem_gb / exa_envs
    per_env_storage = req_storage_tb / exa_envs
    per_env_iops = req_iops / exa_envs

    db_servers_per_env = max(
        1,
        int(
            math.ceil(
                max(
                    per_env_cpu / exa_server_capacity["db_ecpu"] if exa_server_capacity["db_ecpu"] > 0 else 0.0,
                    per_env_mem / exa_server_capacity["db_memory_gb"] if exa_server_capacity["db_memory_gb"] > 0 else 0.0,
                )
            )
        ),
    )
    storage_servers_per_env = max(
        1,
        int(
            math.ceil(
                max(
                    per_env_storage / exa_server_capacity["storage_tb"] if exa_server_capacity["storage_tb"] > 0 else 0.0,
                    per_env_iops / exa_server_capacity["storage_iops"] if exa_server_capacity["storage_iops"] > 0 else 0.0,
                )
            )
        ),
    )
    exa_db_servers_total = db_servers_per_env * exa_envs
    exa_storage_servers_total = storage_servers_per_env * exa_envs
    exa_fit = True
    exa_reason = (
        "Fit in a single environment."
        if exa_envs == 1
        else f"Fit with scale-out: single environment is insufficient; requires {exa_envs} environments."
    )

    if base_fit:
        tier, pos, confidence = "Base Database", 1, 0.86
    elif xs_fit:
        tier, pos, confidence = "Exascale", 2, 0.80
    else:
        tier, pos, confidence = "Exadata Dedicated", 3, 0.78

    inventory_rows = [
        {
            "platform": "Base Database",
            "fit": "Fit" if base_fit else "Not Fit",
            "inventory": f"{total_instances} deployments | {base_db_servers} DB servers",
            "reason": base_reason,
        },
        {
            "platform": "Exascale",
            "fit": "Fit" if xs_fit else "Not Fit",
            "inventory": f"{total_cohorts} deployments (1 per cohort)",
            "reason": xs_reason,
        },
        {
            "platform": "Exadata Dedicated",
            "fit": "Fit" if exa_fit else "Not Fit",
            "inventory": (
                f"{exa_envs} environment(s) | {exa_db_servers_total} DB servers | "
                f"{exa_storage_servers_total} Storage servers"
            ),
            "reason": exa_reason,
        },
    ]

    footnotes = [
        "[1] Base DB thresholds used by script: 256 ECPU, 512 GB memory, 80 TB storage, 8k IOPS per instance.",
        "[2] Exascale thresholds used by script: 200 ECPU and 100 TB per cohort deployment.",
        "[3] Exadata-family conversion used by script: 1 vCPU = 2 ECPU for Exascale and Exadata Dedicated sizing.",
        "[4] Exadata Dedicated memory policy used by script: 70% of one DB server memory (1390 GB) => 973 GB per cohort limit.",
        "[5] Exadata server sizing assumptions for inventory: ~380 ECPU and ~973 GB memory policy per DB server (70% of 1390 GB); ~300 TB and ~15M IOPS per Storage server.",
    ]

    return {
        "tier": tier,
        "pos": pos,
        "confidence": confidence,
        "cohort_matrix": cohort_matrix,
        "inventory_rows": inventory_rows,
        "base_failures": base_failures,
        "xs_failures": xs_failures,
        "xs_cohort_details": xs_cohort_details,
        "exadata_envs": exa_envs,
        "base_limits": base_limits,
        "xs_limits": xs_limits,
        "exa_env_limits": exa_env_limits,
        "exa_one_db_server_memory_gb": exa_one_db_server_memory_gb,
        "exa_cohort_memory_limit_gb": exa_cohort_memory_limit_gb,
        "exa_server_capacity": exa_server_capacity,
        "base_db_servers": base_db_servers,
        "exadata_db_servers_total": exa_db_servers_total,
        "exadata_storage_servers_total": exa_storage_servers_total,
        "required_vcpu_total": req_vcpu,
        "required_ecpu_total": req_cpu,
        "required_memory_gb_total": req_mem_gb,
        "required_storage_tb_total": req_storage_tb,
        "required_iops_total": req_iops,
        "footnotes": footnotes,
    }


def infra_progress_values(pos: int, score: float) -> list[float]:
    s = max(0.0, min(100.0, float(score)))
    if pos <= 1:
        p1 = min(100.0, (s / 35.0) * 100.0)
        return [p1, 0.0, 0.0]
    if pos == 2:
        p2 = min(100.0, ((s - 35.0) / 35.0) * 100.0)
        return [100.0, max(0.0, p2), 0.0]
    p3 = min(100.0, ((s - 70.0) / 30.0) * 100.0)
    return [100.0, 100.0, max(0.0, p3)]


def add_infra_progress_chart(slide, values: list[float], x: float, y: float, w: float, h: float):
    data = CategoryChartData()
    data.categories = ["Base Database", "Exascale", "Exadata Dedicated"]
    data.add_series("Infrastructure fit (%)", values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
        data,
    ).chart
    chart.has_legend = False
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 100
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.tick_labels.font.name = "Oracle Sans Tab"
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.name = "Oracle Sans Tab"
    chart.category_axis.tick_labels.font.size = Pt(9)
    plot = chart.plots[0]
    plot.vary_by_categories = True
    try:
        pts = chart.series[0].points
        if len(pts) >= 3:
            pts[0].format.fill.solid()
            pts[0].format.fill.fore_color.rgb = ORACLE_COLORS["accent6"]
            pts[1].format.fill.solid()
            pts[1].format.fill.fore_color.rgb = ORACLE_COLORS["accent4"]
            pts[2].format.fill.solid()
            pts[2].format.fill.fore_color.rgb = ORACLE_COLORS["accent3"]
    except Exception:
        pass


def summary_dba_learnings(
    model: "Model",
    s: dict[str, float],
    total_db: int,
    total_instances: int,
    total_hosts: int,
    total_cohorts: int,
    versions: list[tuple[str, int]],
    cpu_stats: list[tuple[str, dict[str, float]]],
    iops_items: list[tuple[str, float]],
) -> list[str]:
    lines: list[str] = []
    lines.append(
        f"Scope baseline: {total_db} DBs, {total_instances} instances, {total_hosts} hosts across {total_cohorts} cohorts."
    )
    if versions:
        top_ver, top_ver_cnt = versions[0]
        lines.append(f"Most common version is {top_ver} ({top_ver_cnt}/{max(total_db,1)} DBs, {pct(top_ver_cnt, max(total_db,1)):.1f}%).")

    used = float(s.get("used_storage_gb") or 0.0)
    alloc = float(s.get("allocated_storage_gb") or 0.0)
    util = pct(used, alloc)
    if alloc > 0:
        lines.append(f"Storage utilization is {util:.1f}% ({fmt_num(used,1)} GB used of {fmt_num(alloc,1)} GB allocated).")
        if util >= 80.0:
            lines.append("Storage headroom is limited; prioritize growth planning and tiering before expansion events.")

    mem_total = float(s.get("memory_gb_total") or 0.0)
    mem_ded = float(s.get("sga_total") or 0.0) + float(s.get("pga_total") or 0.0)
    mem_share = pct(mem_ded, mem_total)
    if mem_total > 0:
        lines.append(f"Instance-dedicated memory (SGA+PGA) is {mem_share:.1f}% of host memory in scope.")
        if mem_share < 45.0:
            lines.append("Memory appears over-provisioned for current workload; consider right-sizing or consolidation.")

    if iops_items:
        top_iops = max(iops_items, key=lambda kv: kv[1])
        lines.append(f"Highest cohort IOPS is {top_iops[0]} at {fmt_num(top_iops[1],0)} IOPS.")
        if top_iops[1] >= 50000:
            lines.append("Sustained high IOPS suggests evaluating Exadata or equivalent high-throughput architecture.")

    if cpu_stats:
        top_cpu = max(cpu_stats, key=lambda kv: kv[1].get("p95", 0.0))
        lines.append(f"Highest CPU p95 is in {top_cpu[0]} at {fmt_num(top_cpu[1].get('p95',0.0),2)}.")

    return lines[:5]


def cohort_dba_learnings(
    model: "Model",
    cohort: str,
    rs: list["Row"],
    sm: dict[str, float],
    versions: list[tuple[str, int]],
    inst_stats: list[tuple[str, dict[str, float]]],
    mem_items_c: list[tuple[str, float]],
    iops_items_c: list[tuple[str, float]],
) -> list[str]:
    total_db = int(sm["db_count"])
    total_hosts = int(sm["host_count"])
    total_instances = int(sm["instance_count"])
    lines: list[str] = [
        f"Cohort baseline: {total_db} DBs, {total_instances} instances, {total_hosts} hosts."
    ]
    if versions:
        top_ver, top_ver_cnt = versions[0]
        lines.append(f"Most common DB version is {top_ver} ({top_ver_cnt}/{max(total_db,1)} DBs, {pct(top_ver_cnt, max(total_db,1)):.1f}%).")

    mem_total = float(sm.get("memory_gb_total") or 0.0)
    mem_ded = float(sm.get("sga_total") or 0.0) + float(sm.get("pga_total") or 0.0)
    mem_share = pct(mem_ded, mem_total)
    if mem_total > 0:
        lines.append(f"SGA+PGA is {mem_share:.1f}% of host memory in this cohort.")
        if mem_share < 45.0:
            lines.append("Cohort indicates memory over-allocation versus database memory footprint.")

    if inst_stats:
        top_cpu = max(inst_stats, key=lambda kv: kv[1].get("p95", 0.0))
        lines.append(f"Highest instance CPU p95: {short_label(top_cpu[0],20)} at {fmt_num(top_cpu[1].get('p95',0.0),2)}.")

    if iops_items_c:
        top_iops_db = max(iops_items_c, key=lambda kv: kv[1])
        lines.append(f"Highest DB IOPS: {short_label(top_iops_db[0],20)} at {fmt_num(top_iops_db[1],0)}.")
        if top_iops_db[1] >= 30000:
            lines.append("This cohort has high IO pressure and is a potential candidate for Exadata consolidation.")

    if mem_items_c:
        top_mem_db = max(mem_items_c, key=lambda kv: kv[1])
        lines.append(f"Highest DB memory p95: {short_label(top_mem_db[0],20)} at {fmt_num(top_mem_db[1],1)} GB.")

    return lines[:5]


def instance_dba_comment(model: "Model", cohort: str, db: str, inst: str, row: "Row" | None) -> str:
    cpu_stats = model.db_metric_stats(cohort, db, "DB vCPU") or {}
    iops_stats = model.db_metric_stats(cohort, db, "DB IOPS") or {}
    mem_stats = model.db_metric_stats(cohort, db, "DB Memory (MB)") or {}
    mem_total = float(row.mem_gb if row else 0.0)
    mem_ded = float((row.sga_gb if row else 0.0) + (row.pga_gb if row else 0.0))
    mem_share = pct(mem_ded, mem_total)
    cpu_p95 = float(cpu_stats.get("p95") or 0.0)
    iops_p95 = float(iops_stats.get("p95") or 0.0)
    mem_p95_gb = float(mem_stats.get("p95") or 0.0) / 1024.0
    recommendations: list[str] = []

    if mem_total > 0 and mem_share < 40.0:
        recommendations.append(f"Dedicated DB memory is {mem_share:.1f}% of host memory; consider right-sizing.")
    if iops_p95 >= 30000 or cpu_p95 >= 24.0:
        recommendations.append("Observed IO/CPU pressure is elevated; evaluate Exadata for performance and consolidation.")
    if mem_p95_gb > 0 and mem_total > 0 and mem_p95_gb < (0.5 * mem_total):
        recommendations.append("DB memory demand is materially below host memory; review consolidation opportunities.")
    if not recommendations:
        recommendations.append("CPU, memory and IO indicators are within expected range for this instance scope.")

    return short_label(f"DBA assessment: {recommendations[0]}", 120)


@dataclass
class Row:
    cohort: str
    db: str
    instance: str
    host: str
    init_cpu_count: float
    logical_cpu_count: float
    mem_gb: float
    sga_gb: float
    pga_gb: float


class Model:
    def __init__(self, raw: dict[str, Any], ui_state: dict[str, Any] | None = None):
        self.raw = raw
        self.ui_state = ui_state or {}
        self.rows: list[Row] = self._normalize_rows()
        self.cpu_by_cohort = self._parse_cpu_series()
        self.snap_time_by_cohort = self._parse_snap_time_axis()
        self.db_stats = self._normalize_db_stats()
        self.db_versions = self._normalize_db_versions()
        self.db_names = self._normalize_db_names()
        self.db_meta = self._normalize_db_meta()
        self.db_params = self._normalize_db_params()
        self.cohort_rollups = self._normalize_cohort_rollups()
        self.selection = self._resolve_selection()

    def _normalize_rows(self) -> list[Row]:
        out: list[Row] = []
        for r in self.raw.get("instances", []) or []:
            cohort = r.get("cdb_cohort") or r.get("db_cohort") or "UNASSIGNED"
            db = r.get("cdb") or r.get("WHICH_DB") or "UNKNOWN_DB"
            inst = r.get("db") or f"{db}${r.get('host') or 'UNKNOWN_HOST'}"
            out.append(
                Row(
                    cohort=cohort,
                    db=db,
                    instance=inst,
                    host=r.get("host") or "UNKNOWN_HOST",
                    init_cpu_count=to_num(r.get("init_cpu_count")),
                    logical_cpu_count=to_num(r.get("logical_cpu_count")),
                    mem_gb=to_num(r.get("mem_gb")),
                    sga_gb=to_num(r.get("sga_size_gb")),
                    pga_gb=to_num(r.get("pga_size_gb")),
                )
            )
        return out

    def _parse_cpu_series(self) -> dict[str, tuple[list[str], list[float]]]:
        out: dict[str, tuple[list[str], list[float]]] = {}
        for p in self.raw.get("plots_html", []) or []:
            name = p.get("file_name") or ""
            if not name.endswith("_DB_vCPU_unadjusted.html"):
                continue
            cohort = name.replace("_DB_vCPU_unadjusted.html", "")
            parsed = parse_timeseries_from_plot_html(p.get("file_content") or "")
            if parsed:
                out[cohort] = parsed
                c2 = p.get("cohort")
                if c2 and c2 not in out:
                    out[c2] = parsed
        return out

    def _parse_snap_time_axis(self) -> dict[str, list[str]]:
        out: dict[str, list[str]] = {}
        for p in self.raw.get("plots_html", []) or []:
            name = p.get("file_name") or ""
            if not name.endswith("_awr_snap_coverage.html"):
                continue
            cohort = name.replace("_awr_snap_coverage.html", "")
            xs = parse_time_axis_from_plot_html(p.get("file_content") or "")
            if xs:
                out[cohort] = xs
                c2 = p.get("cohort")
                if c2 and c2 not in out:
                    out[c2] = xs
        return out

    def _normalize_db_stats(self) -> list[dict[str, Any]]:
        out = []
        for r in self.raw.get("database_statistics", []) or []:
            out.append(
                {
                    "cohort": r.get("cdb_cohort") or "UNASSIGNED",
                    "db": r.get("cdb") or "UNKNOWN_DB",
                    "metric": r.get("metric") or "",
                    "min": to_num(r.get("min")),
                    "p30": to_num(r.get("p30")),
                    "p50": to_num(r.get("p50")),
                    "p70": to_num(r.get("p70")),
                    "p95": to_num(r.get("p95")),
                    "p99": to_num(r.get("p99")),
                    "max": to_num(r.get("max")),
                    "summary_of": r.get("summary_of") or "",
                }
            )
        return out

    def _normalize_db_versions(self) -> dict[str, str]:
        out: dict[str, str] = {}
        for d in self.raw.get("databases", []) or []:
            db = d.get("cdb") or d.get("WHICH_DB")
            if not db:
                continue
            out[db] = d.get("cdb_version") or d.get("db_version") or "Unknown"
        return out

    def _normalize_db_names(self) -> dict[str, str]:
        out: dict[str, str] = {}
        for r in self.raw.get("properties_database", []) or []:
            db = r.get("WHICH_DB") or r.get("cdb")
            if not db:
                continue
            nm = r.get("DB_NAME")
            if nm:
                out[db] = str(nm)
        return out

    def _normalize_db_meta(self) -> dict[str, dict[str, Any]]:
        out: dict[str, dict[str, Any]] = {}
        for d in self.raw.get("databases", []) or []:
            db = d.get("cdb") or d.get("WHICH_DB")
            if not db:
                continue
            out[db] = {
                "version": d.get("cdb_version") or d.get("db_version") or "Unknown",
                "instance_count": int(to_num(d.get("cdb_instance_count"), 0)),
                "clustered": d.get("clustered") or "",
            }
        return out

    def _normalize_db_params(self) -> dict[str, dict[str, str]]:
        out: dict[str, dict[str, str]] = defaultdict(dict)
        for r in self.raw.get("database_parameters", []) or []:
            db = r.get("cdb")
            pn = str(r.get("PARAMETER_NAME") or "").strip().lower()
            if not db or not pn:
                continue
            if pn in {"enable_pluggable_database", "target_pdbs", "max_pdbs", "noncdb_compatible"}:
                out[db][pn] = str(r.get("VALUE") or "").strip()
        return dict(out)

    def _normalize_cohort_rollups(self) -> dict[str, dict[str, float]]:
        out: dict[str, dict[str, float]] = {}
        for r in self.raw.get("cohort_rollups", []) or []:
            c = r.get("Name")
            if not c:
                continue
            out[c] = {
                "allocated": to_num(r.get("Allocated Storage (GB)")),
                "used": to_num(r.get("Used Storage (GB)")),
                "db_iops": to_num(r.get("DB IOPS")),
            }
        return out

    def _resolve_selection(self) -> dict[str, set[str]]:
        rows = self.rows
        all_cohorts = {r.cohort for r in rows}
        all_dbs = {r.db for r in rows}
        all_inst = {r.instance for r in rows}
        ui = self.ui_state.get("selection", {}).get("ppt", {}) if isinstance(self.ui_state, dict) else {}
        cohorts = set(ui.get("cohorts", []) or all_cohorts)
        dbs = set(ui.get("dbs", []) or all_dbs)
        inst = set(ui.get("instances", []) or all_inst)
        return {"cohorts": cohorts, "dbs": dbs, "instances": inst}

    def selected_rows(self) -> list[Row]:
        sel = self.selection
        return [r for r in self.rows if r.cohort in sel["cohorts"] and r.db in sel["dbs"] and r.instance in sel["instances"]]

    def cohorts(self) -> list[str]:
        return sorted({r.cohort for r in self.selected_rows()})

    def rows_for_cohort(self, cohort: str) -> list[Row]:
        return [r for r in self.selected_rows() if r.cohort == cohort]

    def display_db_name(self, db: str) -> str:
        nm = self.db_names.get(db)
        if nm:
            return nm
        token = str(db or "UNKNOWN_DB")
        if token.startswith("db_"):
            token = token[3:]
        token = token.split("_")[0] if "_" in token else token
        return token or "UNKNOWN_DB"

    def instance_label(self, inst: str) -> str:
        token = str(inst or "")
        m = re.match(r"^db_([^_$]+)", token, re.IGNORECASE)
        if m and m.group(1):
            return m.group(1)
        parts = token.split("$")
        if len(parts) >= 3:
            return f"{parts[-2]}${parts[-1]}"
        return token or "UNKNOWN_INSTANCE"

    def db_type_and_pdb_count(self, db: str) -> tuple[str, str]:
        params = self.db_params.get(db, {})
        epd = params.get("enable_pluggable_database", "").upper()
        target_pdbs = params.get("target_pdbs", "")
        db_type = "Non-CDB"
        if epd == "TRUE":
            db_type = "CDB"
        elif epd == "FALSE":
            db_type = "Non-CDB"

        if db_type == "CDB":
            try:
                return db_type, str(max(0, int(float(target_pdbs))))
            except Exception:
                return db_type, "N/A"
        return db_type, "0"

    def db_version(self, db: str) -> str:
        meta = self.db_meta.get(db, {})
        v = meta.get("version") or self.db_versions.get(db) or "Unknown"
        return str(v)

    def db_is_rac(self, db: str) -> bool:
        meta = self.db_meta.get(db, {})
        clustered = str(meta.get("clustered") or "").upper()
        return "CLUSTERED" in clustered and "NON" not in clustered

    def summary(self, rows: list[Row] | None = None) -> dict[str, float]:
        rs = rows if rows is not None else self.selected_rows()
        cohorts = {r.cohort for r in rs}
        allocated = sum((self.cohort_rollups.get(c, {}).get("allocated") or 0.0) for c in cohorts)
        used = sum((self.cohort_rollups.get(c, {}).get("used") or 0.0) for c in cohorts)
        return {
            "db_count": len({r.db for r in rs}),
            "host_count": len({r.host for r in rs}),
            "instance_count": len(rs),
            "vcpu_total": sum((r.init_cpu_count or r.logical_cpu_count) for r in rs),
            "memory_gb_total": sum(r.mem_gb for r in rs),
            "sga_total": sum(r.sga_gb for r in rs),
            "pga_total": sum(r.pga_gb for r in rs),
            "allocated_storage_gb": allocated,
            "used_storage_gb": used,
        }

    def by_cohort_table(self) -> list[dict[str, Any]]:
        groups: dict[str, dict[str, Any]] = {}
        for r in self.selected_rows():
            g = groups.setdefault(
                r.cohort,
                {"cohort": r.cohort, "dbs": set(), "hosts": set(), "instances": 0, "vcpu": 0.0, "mem": 0.0},
            )
            g["dbs"].add(r.db)
            g["hosts"].add(r.host)
            g["instances"] += 1
            g["vcpu"] += r.init_cpu_count or r.logical_cpu_count
            g["mem"] += r.mem_gb
        out = []
        for g in groups.values():
            out.append(
                {
                    "cohort": g["cohort"],
                    "dbs": len(g["dbs"]),
                    "hosts": len(g["hosts"]),
                    "instances": g["instances"],
                    "vcpu": g["vcpu"],
                    "mem": g["mem"],
                }
            )
        return sorted(out, key=lambda x: x["cohort"])

    def version_counts(self) -> list[tuple[str, int]]:
        cnt: dict[str, int] = defaultdict(int)
        selected_db = sorted({r.db for r in self.selected_rows()})
        for db in selected_db:
            cnt[self.db_versions.get(db, "Unknown")] += 1
        return sorted(cnt.items(), key=lambda kv: (-kv[1], kv[0]))

    def cpu_series_global(self) -> list[tuple[str, list[str], list[float]]]:
        out = []
        for c in self.cohorts():
            s = self.cpu_by_cohort.get(c)
            if s:
                out.append((c, s[0], s[1]))
        return out

    def cpu_stats_by_cohort(self) -> list[tuple[str, dict[str, float]]]:
        out: list[tuple[str, dict[str, float]]] = []
        for c, _xs, ys in self.cpu_series_global():
            st = stats_from_values(ys)
            if st:
                out.append((c, st))
        return out

    def memory_by_cohort(self) -> list[tuple[str, float]]:
        rows = self.by_cohort_table()
        return sorted([(r["cohort"], float(r["mem"])) for r in rows if float(r["mem"]) > 0], key=lambda kv: kv[0])

    def iops_by_cohort(self) -> list[tuple[str, float]]:
        out: list[tuple[str, float]] = []
        for c in self.cohorts():
            v = float((self.cohort_rollups.get(c, {}) or {}).get("db_iops") or 0.0)
            if v > 0:
                out.append((c, v))
        return sorted(out, key=lambda kv: kv[0])

    def cpu_series_instances_in_cohort(self, cohort: str) -> list[tuple[str, list[str], list[float]]]:
        base = self.cpu_by_cohort.get(cohort)
        if not base:
            return []
        xs, ys = base
        rows = self.rows_for_cohort(cohort)
        if not rows:
            return []
        total = sum(r.init_cpu_count for r in rows)
        share_default = 1.0 / len(rows)
        out = []
        for r in rows:
            share = (r.init_cpu_count / total) if total > 0 else share_default
            out.append((r.instance, xs, [v * share for v in ys]))
        return out

    def mem_by_db(self, cohort: str, db_scope: set[str] | None = None) -> list[tuple[str, float]]:
        rows = [
            x
            for x in self.db_stats
            if x["cohort"] == cohort
            and x["metric"] == "DB Memory (MB)"
            and (not x["summary_of"] or x["summary_of"] == "cdb")
            and (db_scope is None or x["db"] in db_scope)
        ]
        groups: dict[str, list[dict[str, Any]]] = defaultdict(list)
        for r in rows:
            groups[r["db"]].append(r)
        out = []
        for db, arr in groups.items():
            p95 = sum(x["p95"] for x in arr) / max(len(arr), 1) / 1024.0
            mx = sum(x["max"] for x in arr) / max(len(arr), 1) / 1024.0
            out.append((db, p95 if p95 > 0 else mx))
        return sorted(out, key=lambda kv: (-kv[1], kv[0]))

    def metric_by_db(self, cohort: str, metric: str, db_scope: set[str] | None = None, scale: float = 1.0) -> list[tuple[str, float]]:
        rows = [
            x
            for x in self.db_stats
            if x["cohort"] == cohort
            and x["metric"] == metric
            and (not x["summary_of"] or x["summary_of"] == "cdb")
            and (db_scope is None or x["db"] in db_scope)
        ]
        groups: dict[str, list[dict[str, Any]]] = defaultdict(list)
        for r in rows:
            groups[r["db"]].append(r)
        out = []
        for db, arr in groups.items():
            p95 = (sum(x["p95"] for x in arr) / max(len(arr), 1)) * scale
            mx = (sum(x["max"] for x in arr) / max(len(arr), 1)) * scale
            out.append((db, p95 if p95 > 0 else mx))
        return sorted(out, key=lambda kv: kv[0])

    def db_metric_stats(self, cohort: str, db: str, metric: str) -> dict[str, float] | None:
        rows = [
            x
            for x in self.db_stats
            if x["cohort"] == cohort
            and x["db"] == db
            and x["metric"] == metric
            and (not x["summary_of"] or x["summary_of"] == "cdb")
        ]
        if not rows:
            rows = [x for x in self.db_stats if x["cohort"] == cohort and x["db"] == db and x["metric"] == metric]
        if not rows:
            return None
        keys = ["min", "p30", "p50", "p70", "p95", "p99", "max"]
        out: dict[str, float] = {}
        for k in keys:
            out[k] = sum(to_num(r.get(k)) for r in rows) / max(len(rows), 1)
        return out

    def _cohort_time_range(self, cohort: str) -> tuple[datetime, datetime] | None:
        keys = list(self.cpu_by_cohort.keys())
        snap_keys = list(self.snap_time_by_cohort.keys())

        def resolve_cpu_for_cohort(name: str):
            if name in self.cpu_by_cohort:
                return self.cpu_by_cohort[name]
            up = name.upper()
            for k in keys:
                if str(k).upper() == up:
                    return self.cpu_by_cohort[k]
            for k in keys:
                ks = str(k).upper()
                if ks.startswith(up) or up.startswith(ks):
                    return self.cpu_by_cohort[k]
            return None

        def resolve_snap_for_cohort(name: str):
            if name in self.snap_time_by_cohort:
                return self.snap_time_by_cohort[name]
            up = name.upper()
            for k in snap_keys:
                if str(k).upper() == up:
                    return self.snap_time_by_cohort[k]
            for k in snap_keys:
                ks = str(k).upper()
                if ks.startswith(up) or up.startswith(ks):
                    return self.snap_time_by_cohort[k]
            return None

        def parse_dt(xs: list[str]) -> list[datetime]:
            ts: list[datetime] = []
            for x in xs:
                txt = str(x).replace("Z", "")
                try:
                    ts.append(datetime.fromisoformat(txt))
                except Exception:
                    continue
            return ts

        s = resolve_cpu_for_cohort(cohort)
        ts = parse_dt(s[0]) if s else []
        if not ts:
            xs = resolve_snap_for_cohort(cohort) or []
            ts = parse_dt(xs)
        if not ts:
            return None
        return (min(ts), max(ts))

    def freshness_ranges_by_instance_in_cohort(self, cohort: str) -> list[tuple[str, datetime, datetime]]:
        span = self._cohort_time_range(cohort)
        if not span:
            return []
        start_dt, end_dt = span
        instances = sorted({r.instance for r in self.rows_for_cohort(cohort)})
        return [(inst, start_dt, end_dt) for inst in instances]

    def freshness_ranges(self) -> list[tuple[str, datetime, datetime]]:
        out: list[tuple[str, datetime, datetime]] = []

        for c in self.cohorts():
            span = self._cohort_time_range(c)
            if not span:
                continue
            out.append((c, span[0], span[1]))
        # Fallback: if cohort matching fails, still provide global timeline by available keys.
        if not out:
            for k, s in list(self.cpu_by_cohort.items())[:8]:
                ts = parse_dt(s[0])
                if ts:
                    out.append((str(k), min(ts), max(ts)))
            if not out:
                for k, xs in list(self.snap_time_by_cohort.items())[:8]:
                    ts = parse_dt(xs)
                    if ts:
                        out.append((str(k), min(ts), max(ts)))
        out.sort(key=lambda t: t[0])
        return out


def find_layout(prs: Presentation, preferred_names: list[str], fallback_light: bool = True):
    names = [n.lower() for n in preferred_names]
    for lay in prs.slide_layouts:
        n = (lay.name or "").lower()
        if any(p in n for p in names):
            return lay
    if fallback_light:
        for lay in prs.slide_layouts:
            n = (lay.name or "").lower()
            if "light" in n and "dark" not in n:
                return lay
    return prs.slide_layouts[0]


def presentation_from_template(path: Path) -> Presentation:
    # python-pptx rejects .potx content-type directly; create a temporary
    # .pptx clone with presentation main content-type.
    if path.suffix.lower() != ".potx":
        return Presentation(str(path))
    with zipfile.ZipFile(path, "r") as zin:
        tmpdir = Path(tempfile.mkdtemp(prefix="potx_as_pptx_"))
        zin.extractall(tmpdir)
    ct = tmpdir / "[Content_Types].xml"
    txt = ct.read_text(encoding="utf-8")
    txt = txt.replace(
        "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
    )
    ct.write_text(txt, encoding="utf-8")
    out = tmpdir / "template_as_presentation.pptx"
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
        for p in sorted(tmpdir.rglob("*")):
            if p == out or p.is_dir():
                continue
            zout.write(p, p.relative_to(tmpdir))
    return Presentation(str(out))


def clear_existing_slides(prs: Presentation) -> None:
    """
    Remove any slides already present in the template so export contains only
    generated AWR slides.
    """
    slide_ids = list(prs.slides._sldIdLst)  # pylint: disable=protected-access
    for sld_id in slide_ids:
        r_id = sld_id.rId
        prs.part.drop_rel(r_id)
        prs.slides._sldIdLst.remove(sld_id)  # pylint: disable=protected-access


def set_title_and_subtitle(slide, title: str, subtitle: str = ""):
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
        tf = slide.shapes.title.text_frame
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.name = "Oracle Sans Tab"
                r.font.color.rgb = ORACLE_COLORS["text"]
    # First body placeholder as subtitle when available
    body_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.type == 2:  # BODY
            body_ph = ph
            break
    if body_ph is not None and subtitle:
        body_ph.text = subtitle
        tf = body_ph.text_frame
        for p in tf.paragraphs:
            for r in p.runs:
                r.font.name = "Oracle Sans Tab"
                r.font.color.rgb = ORACLE_COLORS["muted"]


def add_kpi_boxes(slide, kpis: list[tuple[str, str]]):
    # Content region in white template is roughly 0.84..12.5 x 1.75..6.68
    cols = 4
    x0 = 0.9
    y0 = 1.75
    w = 2.85
    h = 0.72
    gx = 0.18
    gy = 0.1
    for i, (k, v) in enumerate(kpis[:8]):
        c = i % cols
        r = i // cols
        x = x0 + c * (w + gx)
        y = y0 + r * (h + gy)
        shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))  # MSO_AUTO_SHAPE_TYPE.RECTANGLE
        shp.fill.solid()
        shp.fill.fore_color.rgb = ORACLE_COLORS["panel"]
        shp.line.color.rgb = ORACLE_COLORS["line"]
        tf = shp.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = k
        p1.font.name = "Oracle Sans Tab"
        p1.font.size = Pt(9)
        p1.font.color.rgb = ORACLE_COLORS["muted"]
        p2 = tf.add_paragraph()
        p2.text = v
        p2.font.name = "Oracle Sans Tab"
        p2.font.bold = True
        p2.font.size = Pt(13)
        p2.font.color.rgb = ORACLE_COLORS["text"]


def add_exec_cards(slide, cards: list[dict[str, str]]):
    x0 = 0.8
    y0 = 1.55
    count = min(len(cards), 5)
    gap = 0.14
    usable_w = 12.0
    w = (usable_w - (gap * (count - 1))) / max(count, 1)
    h = 1.35
    for i, c in enumerate(cards[:5]):
        x = x0 + i * (w + gap)
        box = slide.shapes.add_shape(1, Inches(x), Inches(y0), Inches(w), Inches(h))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFC)
        box.line.color.rgb = ORACLE_COLORS["line"]
        # Top accent strip
        strip = slide.shapes.add_shape(1, Inches(x), Inches(y0), Inches(w), Inches(0.12))
        strip.fill.solid()
        strip.fill.fore_color.rgb = c.get("color", ORACLE_COLORS["accent1"])
        strip.line.fill.background()
        tf = box.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        p0.text = c.get("label", "")
        p0.font.name = "Oracle Sans Tab"
        p0.font.bold = True
        p0.font.size = Pt(10)
        p0.font.color.rgb = ORACLE_COLORS["muted"]
        p0.alignment = PP_ALIGN.LEFT
        p1 = tf.add_paragraph()
        p1.text = c.get("value", "")
        p1.font.name = "Oracle Sans Tab"
        p1.font.bold = True
        p1.font.size = Pt(22)
        p1.font.color.rgb = c.get("color", ORACLE_COLORS["accent1"])
        p1.alignment = PP_ALIGN.LEFT
        p2 = tf.add_paragraph()
        p2.text = c.get("sub", "")
        p2.font.name = "Oracle Sans Tab"
        p2.font.size = Pt(9)
        p2.font.color.rgb = ORACLE_COLORS["muted"]
        p2.alignment = PP_ALIGN.LEFT


def add_big_versions_card(slide, title: str, lines: list[str], x: float, y: float, w: float, h: float):
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFC)
    box.line.color.rgb = ORACLE_COLORS["line"]
    strip = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.12))
    strip.fill.solid()
    strip.fill.fore_color.rgb = ORACLE_COLORS["accent2"]
    strip.line.fill.background()
    tf = box.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.name = "Oracle Sans Tab"
    p0.font.bold = True
    p0.font.size = Pt(18)
    p0.font.color.rgb = ORACLE_COLORS["text"]
    p0.alignment = PP_ALIGN.LEFT
    for line in lines[:5]:
        p = tf.add_paragraph()
        p.text = line
        p.font.name = "Oracle Sans Tab"
        p.font.size = Pt(14)
        p.font.color.rgb = ORACLE_COLORS["muted"]
        p.alignment = PP_ALIGN.LEFT


def add_panel(
    slide,
    title: str,
    lines: list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    dark: bool = False,
    body_font_size: int = 12,
):
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x23, 0x2B, 0x62) if dark else RGBColor(0xF7, 0xF9, 0xFC)
    box.line.color.rgb = ORACLE_COLORS["line"]
    if not dark:
        left = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.12), Inches(h))
        left.fill.solid()
        left.fill.fore_color.rgb = ORACLE_COLORS["accent1"]
        left.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Cm(0.5)
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.name = "Oracle Sans Tab"
    p0.font.bold = True
    p0.font.size = Pt(18)
    p0.font.color.rgb = RGBColor(0xE7, 0xEE, 0xFF) if dark else ORACLE_COLORS["accent1"]
    p0.alignment = PP_ALIGN.LEFT
    for ln in lines[:8]:
        p = tf.add_paragraph()
        p.text = ln
        p.font.name = "Oracle Sans Tab"
        p.font.size = Pt(body_font_size)
        p.font.color.rgb = RGBColor(0xE7, 0xEE, 0xFF) if dark else ORACLE_COLORS["text"]
        p.level = 0
        p.alignment = PP_ALIGN.LEFT


def add_footer_note(slide, text: str, x: float, y: float, w: float):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Oracle Sans Tab"
    p.font.size = Pt(7)
    p.font.color.rgb = ORACLE_COLORS["muted"]
    p.alignment = PP_ALIGN.LEFT


def add_section_label(slide, text: str, x: float, y: float):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(3.8), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Oracle Sans Tab"
    p.font.bold = True
    p.font.size = Pt(10)
    p.font.color.rgb = ORACLE_COLORS["muted"]


def add_insight_callout(slide, title: str, lines: list[str], x: float, y: float, w: float, h: float):
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = ORACLE_COLORS["panel"]
    shp.line.color.rgb = ORACLE_COLORS["line"]
    tf = shp.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.name = "Oracle Sans Tab"
    p0.font.bold = True
    p0.font.size = Pt(10)
    p0.font.color.rgb = ORACLE_COLORS["text"]
    for line in lines[:3]:
        p = tf.add_paragraph()
        p.text = f"- {line}"
        p.font.name = "Oracle Sans Tab"
        p.font.size = Pt(9)
        p.font.color.rgb = ORACLE_COLORS["muted"]


def add_notes(slide, lines: list[str], x: float, y: float, w: float, h: float):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(lines[:3]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"- {line}"
        p.font.name = "Oracle Sans Tab"
        p.font.size = Pt(9)
        p.font.color.rgb = ORACLE_COLORS["muted"]


def add_table(
    slide,
    headers: list[str],
    rows: list[list[str]],
    x: float,
    y: float,
    w: float,
    h: float,
    numeric_cols: set[int] | None = None,
    fit_cols: set[int] | None = None,
):
    numeric_cols = numeric_cols or set()
    fit_cols = fit_cols or set()
    t = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for j, htxt in enumerate(headers):
        cell = t.cell(0, j)
        cell.text = htxt
        cell.fill.solid()
        cell.fill.fore_color.rgb = ORACLE_COLORS["accent1"]
        p = cell.text_frame.paragraphs[0]
        p.font.name = "Oracle Sans Tab"
        p.font.bold = True
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.RIGHT if j in numeric_cols else PP_ALIGN.LEFT
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            cell = t.cell(i, j)
            cell.text = v
            txt = str(v or "").strip().lower()
            if j in fit_cols and txt.startswith("fit"):
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xD9, 0xF2, 0xE3)  # light green
            elif j in fit_cols and txt.startswith("no fit"):
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF9, 0xDD, 0xDD)  # light red
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF5, 0xF7) if i % 2 else RGBColor(0xE8, 0xED, 0xF1)
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Oracle Sans Tab"
            p.font.size = Pt(8)
            p.font.color.rgb = ORACLE_COLORS["text"]
            p.alignment = PP_ALIGN.RIGHT if j in numeric_cols else PP_ALIGN.LEFT


def downsample(xs: list[str], ys: list[float], max_points: int = 36) -> tuple[list[str], list[float]]:
    n = min(len(xs), len(ys))
    if n <= max_points:
        return xs[:n], ys[:n]
    out_i = [round(i * (n - 1) / (max_points - 1)) for i in range(max_points)]
    out_i = sorted(set(out_i))
    return [xs[i] for i in out_i], [ys[i] for i in out_i]


def add_line_chart(
    slide,
    series: list[tuple[str, list[str], list[float]]],
    x: float,
    y: float,
    w: float,
    h: float,
    show_legend: bool = True,
    y_min: float | None = None,
    y_max: float | None = None,
    x_tick_font_size: int = 8,
    y_tick_font_size: int = 8,
    chart_kind: str = "line",
    line_width_pt: float = 2.0,
    area_transparency: float = 0.0,
):
    if not series:
        return
    chart_data = CategoryChartData()
    base_x = series[0][1]
    if not base_x:
        return
    labels, _ = downsample(base_x, series[0][2], 36)
    chart_data.categories = [str(t).replace("T", " ")[5:16] for t in labels]
    top = sorted(series, key=lambda s: max(s[2]) if s[2] else 0.0, reverse=True)[:8]
    for name, xs, ys in top:
        sx, sy = downsample(xs, ys, 36)
        # align by downsample count
        m = min(len(labels), len(sx), len(sy))
        chart_data.add_series(short_label(name, 24), sy[:m])
    chart_type = XL_CHART_TYPE.AREA if str(chart_kind).lower() == "area" else XL_CHART_TYPE.LINE
    chart = slide.shapes.add_chart(
        chart_type,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
        chart_data,
    ).chart
    chart.has_title = False
    chart.has_legend = bool(show_legend)
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        try:
            chart.legend.font.size = Pt(10)
        except Exception:
            pass
    if y_min is not None:
        chart.value_axis.minimum_scale = float(y_min)
    if y_max is not None:
        chart.value_axis.maximum_scale = float(y_max)
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.has_minor_gridlines = False
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.has_minor_gridlines = False
    chart.value_axis.tick_labels.font.size = Pt(y_tick_font_size)
    chart.category_axis.tick_labels.font.size = Pt(x_tick_font_size)
    for i, s in enumerate(chart.series):
        color = PALETTE[i % len(PALETTE)]
        s.format.line.color.rgb = color
        s.format.line.width = Pt(line_width_pt)
        if chart_type == XL_CHART_TYPE.AREA:
            try:
                s.format.fill.solid()
                s.format.fill.fore_color.rgb = color
                s.format.fill.transparency = max(0.0, min(1.0, float(area_transparency)))
                # Force alpha in chart-series XML for PowerPoint reliability.
                ser = s._element
                sp_pr = ser.find(qn("c:spPr"))
                if sp_pr is None:
                    sp_pr = OxmlElement("c:spPr")
                    ser.append(sp_pr)
                for node in list(sp_pr.findall(qn("a:solidFill"))):
                    sp_pr.remove(node)
                solid = OxmlElement("a:solidFill")
                srgb = OxmlElement("a:srgbClr")
                srgb.set("val", f"{int(color[0]):02X}{int(color[1]):02X}{int(color[2]):02X}")
                alpha = OxmlElement("a:alpha")
                # OOXML alpha: 100000 = fully opaque, 0 = fully transparent.
                alpha_val = int(round((1.0 - max(0.0, min(1.0, float(area_transparency)))) * 100000))
                alpha.set("val", str(alpha_val))
                srgb.append(alpha)
                solid.append(srgb)
                sp_pr.append(solid)
            except Exception:
                pass


def add_bar_chart(
    slide,
    items: list[tuple[str, float]],
    x: float,
    y: float,
    w: float,
    h: float,
    series_name: str = "Value",
    x_tick_font_size: int = 8,
):
    vals = [(short_label(k, 22), float(v)) for k, v in items if float(v) > 0][:10]
    if not vals:
        return
    data = CategoryChartData()
    data.categories = [k for k, _ in vals]
    data.add_series(series_name, [v for _, v in vals])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
        data,
    ).chart
    chart.has_title = False
    chart.has_legend = False
    chart.value_axis.minimum_scale = 0
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.has_minor_gridlines = False
    chart.category_axis.has_major_gridlines = False
    chart.category_axis.has_minor_gridlines = False
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(x_tick_font_size)
    plot = chart.plots[0]
    plot.vary_by_categories = True


def add_vcpu_boxplot_chart(
    slide,
    stats: dict[str, float] | None,
    x: float,
    y: float,
    w: float,
    h: float,
    title_text: str = "CPU Consumption (Box Plot)",
    value_scale: float = 1.0,
):
    title = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.18))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Oracle Sans Tab"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ORACLE_COLORS["muted"]
    p.alignment = PP_ALIGN.LEFT

    if not stats:
        add_freshness_unavailable_note(slide, x, y + 0.28, w)
        return

    vmin = to_num(stats.get("min")) * value_scale
    p30 = to_num(stats.get("p30")) * value_scale
    p50 = to_num(stats.get("p50")) * value_scale
    p70 = to_num(stats.get("p70")) * value_scale
    vmax = to_num(stats.get("max")) * value_scale
    p95 = to_num(stats.get("p95")) * value_scale
    # Real Python boxplot (matplotlib), inserted as a rendered chart image.
    try:
        os.environ.setdefault("MPLCONFIGDIR", "/tmp/mplconfig")
        import matplotlib

        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        fig_w = max(3.0, w * 1.6)
        fig_h = max(1.2, (h - 0.24) * 1.6)
        fig, ax = plt.subplots(figsize=(fig_w, fig_h), dpi=180)
        bxpstats = [
            {
                "label": "",
                "whislo": vmin,
                "q1": p30,
                "med": p50,
                "q3": p70,
                "whishi": vmax,
                "fliers": [p95] if (p95 < vmin or p95 > vmax) else [],
            }
        ]
        b = ax.bxp(bxpstats, patch_artist=True, showfliers=True, widths=0.45)
        for box in b.get("boxes", []):
            box.set_facecolor("#DCECF0")
            box.set_edgecolor("#04536F")
            box.set_linewidth(1.2)
        for med in b.get("medians", []):
            med.set_color("#C74634")
            med.set_linewidth(1.4)
        for wk in b.get("whiskers", []) + b.get("caps", []):
            wk.set_color("#04536F")
            wk.set_linewidth(1.1)
        for fl in b.get("fliers", []):
            fl.set_marker("D")
            fl.set_markerfacecolor("#6C3F49")
            fl.set_markeredgecolor("#6C3F49")
            fl.set_markersize(4)

        ax.set_xticks([])
        ax.set_ylabel("")
        ax.tick_params(axis="y", labelsize=8, colors="#6B747A")
        for tl in ax.get_yticklabels():
            tl.set_fontname("Oracle Sans Tab")
        ax.grid(axis="y", linestyle="-", linewidth=0.6, color="#D9DEDE", alpha=0.85)
        for sp in ("top", "right"):
            ax.spines[sp].set_visible(False)
        ax.spines["left"].set_color("#D9DEDE")
        ax.spines["bottom"].set_color("#D9DEDE")
        ax.set_facecolor((1, 1, 1, 0))
        fig.patch.set_alpha(0.0)
        fig.tight_layout(pad=0.4)

        png = io.BytesIO()
        fig.savefig(png, format="png", dpi=180, transparent=True)
        plt.close(fig)
        png.seek(0)
        slide.shapes.add_picture(
            png,
            Inches(x),
            Inches(y + 0.22),
            Inches(w),
            Inches(max(0.2, h - 0.24)),
        )
    except Exception:
        # Fully native editable fallback chart with percentile distribution.
        data = CategoryChartData()
        data.categories = ["Min", "P30", "P50", "P70", "P95", "Max"]
        data.add_series("CPU", [vmin, p30, p50, p70, p95, vmax])
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(x),
            Inches(y + 0.22),
            Inches(w),
            Inches(max(0.2, h - 0.24)),
            data,
        ).chart
        chart.has_legend = False
        try:
            chart.value_axis.tick_labels.font.size = Pt(8)
            chart.category_axis.tick_labels.font.size = Pt(8)
        except Exception:
            pass

    # Show factual reference values so users can interpret the OHLC mapping.
    note = slide.shapes.add_textbox(Inches(x), Inches(y + h - 0.16), Inches(w), Inches(0.16))
    ntf = note.text_frame
    ntf.clear()
    np = ntf.paragraphs[0]
    np.text = f"Min {fmt_num(vmin,3)} | P30 {fmt_num(p30,3)} | Median {fmt_num(p50,3)} | P70 {fmt_num(p70,3)} | Max {fmt_num(vmax,3)}"
    np.font.name = "Oracle Sans Tab"
    np.font.size = Pt(8)
    np.font.color.rgb = ORACLE_COLORS["muted"]
    np.alignment = PP_ALIGN.LEFT


def add_memory_compare_chart(slide, total_host_mem_gb: float, instance_used_mem_gb: float, x: float, y: float, w: float, h: float):
    server_mem = max(0.0, float(total_host_mem_gb))
    inst_mem = max(0.0, float(instance_used_mem_gb))

    # Base native chart (column): instance dedicated memory.
    col_data = CategoryChartData()
    col_data.categories = ["Server", "Instance"]
    col_data.add_series("Instance Dedicated Memory", [0.0, inst_mem])
    col_chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
        col_data,
    ).chart
    col_chart.has_legend = False
    col_chart.value_axis.minimum_scale = 0
    col_chart.value_axis.tick_labels.font.size = Pt(8)
    col_chart.category_axis.tick_labels.font.size = Pt(8)
    try:
        s0 = col_chart.series[0]
        s0.points[0].format.fill.solid()
        s0.points[0].format.fill.fore_color.rgb = RGBColor(0xE6, 0xEA, 0xEF)
        s0.points[1].format.fill.solid()
        s0.points[1].format.fill.fore_color.rgb = ORACLE_COLORS["accent3"]
    except Exception:
        pass

    # Overlay native chart (line): server total memory.
    line_data = CategoryChartData()
    line_data.categories = ["Server", "Instance"]
    line_data.add_series("Server Total Memory", [server_mem, server_mem])
    line_chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
        line_data,
    ).chart
    line_chart.has_legend = False
    line_chart.has_title = False
    line_chart.value_axis.minimum_scale = 0
    line_chart.value_axis.maximum_scale = max(server_mem, inst_mem) * 1.1 if max(server_mem, inst_mem) > 0 else 1.0
    line_chart.value_axis.tick_labels.font.size = Pt(8)
    line_chart.category_axis.tick_labels.font.size = Pt(8)
    line_chart.category_axis.visible = False
    try:
        line_chart.value_axis.visible = False
    except Exception:
        pass
    try:
        ls = line_chart.series[0]
        ls.format.line.color.rgb = ORACLE_COLORS["accent1"]
        ls.marker.style = 2  # square
    except Exception:
        pass


def add_hbar_chart(slide, items: list[tuple[str, float]], x: float, y: float, w: float, h: float, series_name: str = "Value"):
    vals = [(short_label(k, 20), float(v)) for k, v in items if float(v) > 0][:8]
    if not vals:
        return
    data = CategoryChartData()
    data.categories = [k for k, _ in vals]
    data.add_series(series_name, [v for _, v in vals])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
        data,
    ).chart
    chart.has_legend = False
    chart.value_axis.minimum_scale = 0
    chart.category_axis.reverse_order = True
    plot = chart.plots[0]
    plot.vary_by_categories = True


def add_pie_chart(slide, items: list[tuple[str, float]], x: float, y: float, w: float, h: float, series_name: str = "Distribution"):
    vals = [(short_label(k, 22), float(v)) for k, v in items if float(v) > 0][:8]
    if not vals:
        return
    data = CategoryChartData()
    data.categories = [k for k, _ in vals]
    data.add_series(series_name, [v for _, v in vals])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
        data,
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT


def add_freshness_timeline_chart(
    slide,
    ranges: list[tuple[str, datetime, datetime]],
    x: float,
    y: float,
    w: float,
    h: float,
    legend_font_size: int = 10,
):
    if not ranges:
        return

    def to_excel_serial(dt: datetime) -> float:
        epoch = datetime(1899, 12, 30)
        delta = dt - epoch
        return delta.days + (delta.seconds + delta.microseconds / 1_000_000.0) / 86400.0

    global_start = min(r[1] for r in ranges)
    global_end = max(r[2] for r in ranges)
    if global_end <= global_start:
        global_end = global_start + timedelta(hours=1)
    min_serial = to_excel_serial(global_start)
    max_serial = to_excel_serial(global_end)
    if max_serial <= min_serial:
        max_serial = min_serial + (1.0 / 24.0)
    xy = XyChartData()
    for idx, (cohort, c_start, c_end) in enumerate(ranges, start=1):
        s = xy.add_series(short_label(cohort, 18))
        x0 = max(to_excel_serial(c_start), min_serial)
        x1 = max(to_excel_serial(c_end), x0 + (1.0 / 1440.0))
        yv = float(idx)
        s.add_data_point(x0, yv)
        s.add_data_point(x1, yv)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
        xy,
    ).chart
    chart.has_legend = False
    chart.value_axis.minimum_scale = 0.0
    chart.value_axis.maximum_scale = float(len(ranges) + 1)
    chart.value_axis.major_unit = 1.0
    chart.value_axis.has_major_gridlines = False
    chart.category_axis.minimum_scale = min_serial
    chart.category_axis.maximum_scale = max_serial
    chart.category_axis.tick_labels.number_format = "mmm/dd"
    chart.category_axis.tick_labels.font.size = Pt(8)
    try:
        chart.value_axis.visible = False
    except Exception:
        chart.value_axis.tick_labels.font.size = Pt(1)
        chart.value_axis.format.line.fill.background()
    for i, s in enumerate(chart.series):
        s.format.line.color.rgb = PALETTE[i % len(PALETTE)]
        s.format.line.width = Pt(2)
        s.smooth = False
        # Label only the last point with the series (cohort/instance) name.
        try:
            last_point = s.points[len(s.points) - 1]
            dl = last_point.data_label
            dl.has_text_frame = True
            dl.font.name = "Oracle Sans Tab"
            dl.font.size = Pt(4)
            dl.font.color.rgb = ORACLE_COLORS["muted"]
            tf = dl.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = short_label(ranges[i][0], 18)
            run.font.name = "Oracle Sans Tab"
            run.font.size = Pt(4)
            run.font.color.rgb = ORACLE_COLORS["muted"]
        except Exception:
            pass

    # Axis label helper so time context is explicit.
    add_footer_note(
        slide,
        f"Timeline range: {global_start.strftime('%Y-%m-%d %H:%M')} to {global_end.strftime('%Y-%m-%d %H:%M')}",
        x=x,
        y=y + h - 0.02,
        w=w,
    )


def add_freshness_unavailable_note(slide, x: float, y: float, w: float):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.26))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "AWR Data Freshness Timeline: no CPU timeline data available in current scope."
    p.font.name = "Oracle Sans Tab"
    p.font.size = Pt(11)
    p.font.color.rgb = ORACLE_COLORS["muted"]
    p.alignment = PP_ALIGN.LEFT


def add_cohort_cpu_boxplot_chart(slide, cohort_stats: list[tuple[str, dict[str, float]]], x: float, y: float, w: float, h: float):
    if not cohort_stats:
        add_freshness_unavailable_note(slide, x, y + 0.04, w)
        return
    try:
        os.environ.setdefault("MPLCONFIGDIR", "/tmp/mplconfig")
        import matplotlib

        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        items = sorted(cohort_stats, key=lambda t: t[0])
        labels = [short_label(c, 10) for c, _ in items]
        bxpstats = []
        for c, st in items:
            bxpstats.append(
                {
                    "label": short_label(c, 10),
                    "whislo": to_num(st.get("min")),
                    "q1": to_num(st.get("p30")),
                    "med": to_num(st.get("p50")),
                    "q3": to_num(st.get("p70")),
                    "whishi": to_num(st.get("max")),
                    "fliers": [],
                }
            )

        fig_w = max(2.4, w * 1.8)
        fig_h = max(1.0, h * 1.9)
        fig, ax = plt.subplots(figsize=(fig_w, fig_h), dpi=180)
        b = ax.bxp(bxpstats, patch_artist=True, showfliers=False, widths=0.55)
        for i, box in enumerate(b.get("boxes", [])):
            box.set_facecolor("#DCECF0")
            box.set_edgecolor("#04536F")
            box.set_linewidth(1.1)
        for med in b.get("medians", []):
            med.set_color("#C74634")
            med.set_linewidth(1.3)
        for wk in b.get("whiskers", []) + b.get("caps", []):
            wk.set_color("#04536F")
            wk.set_linewidth(1.0)

        ax.set_ylabel("")
        ax.set_xticks(range(1, len(labels) + 1))
        ax.set_xticklabels(labels, fontsize=7, fontname="Oracle Sans Tab", color="#6B747A", rotation=0)
        ax.tick_params(axis="y", labelsize=7, colors="#6B747A")
        for tl in ax.get_yticklabels():
            tl.set_fontname("Oracle Sans Tab")
        ax.grid(axis="y", linestyle="-", linewidth=0.5, color="#D9DEDE", alpha=0.85)
        for sp in ("top", "right"):
            ax.spines[sp].set_visible(False)
        ax.spines["left"].set_color("#D9DEDE")
        ax.spines["bottom"].set_color("#D9DEDE")
        ax.set_facecolor((1, 1, 1, 0))
        fig.patch.set_alpha(0.0)
        fig.tight_layout(pad=0.25)

        png = io.BytesIO()
        fig.savefig(png, format="png", dpi=180, transparent=True)
        plt.close(fig)
        png.seek(0)
        slide.shapes.add_picture(png, Inches(x), Inches(y), Inches(w), Inches(h))
    except Exception:
        add_freshness_unavailable_note(slide, x, y + 0.04, w)


def build_plan(model: Model) -> list[dict[str, Any]]:
    out: list[dict[str, Any]] = []
    out.append(
        {
            "id": "executive-summary",
            "type": "executive-summary",
            "title": "Executive summary: INTERNAL ORACLE - Not to share with customer",
            "subtitle": "DBA analytical overview and deployment recommendation",
        }
    )
    out.append(
        {
            "id": "deployment-base",
            "type": "deployment-base",
            "title": "Service Fit Metrics - Base Database",
            "subtitle": "Informative slide: fit parameters, formulas, thresholds, and references",
        }
    )
    out.append(
        {
            "id": "deployment-exascale",
            "type": "deployment-exascale",
            "title": "Service Fit Metrics - Exascale",
            "subtitle": "Informative slide: fit parameters, formulas, thresholds, and references",
        }
    )
    out.append(
        {
            "id": "deployment-exadata",
            "type": "deployment-exadata",
            "title": "Service Fit Metrics - Exadata Dedicated",
            "subtitle": "Informative slide: fit parameters, formulas, thresholds, and references",
        }
    )
    out.append({"id": "title", "type": "title", "title": "AWR Analysis", "subtitle": "Generated report overview"})
    out.append(
        {
            "id": "summary-global",
            "type": "summary",
            "title": "Global Summary",
            "subtitle": "All selected databases and cohorts",
        }
    )
    out.append({"id": "sep-cohorts", "type": "separator", "title": "Cohort Analysis", "subtitle": "Cohort-level drill down"})
    for c in model.cohorts():
        out.append({"id": f"cohort:{c}", "type": "cohort", "cohort": c, "title": f"Cohort: {c}", "subtitle": f"Summary for {c}"})
    out.append({"id": "sep-instances", "type": "separator", "title": "Instance Analysis", "subtitle": "Instance-level drill down"})
    for c in model.cohorts():
        out.append({"id": f"sep-instance-cohort:{c}", "type": "separator", "title": f"Instances - {c}", "subtitle": f"Instances in {c}"})
        rs = sorted(model.rows_for_cohort(c), key=lambda r: r.instance)
        for r in rs:
            db_name = model.display_db_name(r.db)
            inst_label = model.instance_label(r.instance)
            out.append(
                {
                    "id": f"instance:{c}:{r.instance}",
                    "type": "instance",
                    "cohort": c,
                    "instance": r.instance,
                    "db": r.db,
                    "title": f"Instance: {inst_label}",
                    "subtitle": f"{db_name} on {r.host}",
                }
            )
    out.append({"id": "closing", "type": "closing", "title": "Thank You", "subtitle": "AWR analysis completed"})
    selected = set(model.ui_state.get("pptSlidesSelected") or []) if isinstance(model.ui_state, dict) else set()
    if selected:
        out = [x for x in out if x.get("id") in selected]
    return out


def export_ppt(input_path: Path, template_path: Path, output_path: Path):
    payload = json.loads(input_path.read_text(encoding="utf-8"))
    if isinstance(payload, dict) and "raw_data" in payload:
        raw = payload.get("raw_data") or {}
        ui = payload.get("ui_state") or {}
    else:
        raw = payload
        ui = {}
    model = Model(raw, ui)
    ai_payload = get_ai_app_payload(ui)
    report_meta = get_report_meta(payload if isinstance(payload, dict) else {}, ui)

    prs = presentation_from_template(template_path)
    # Prefer explicit layout mapping from template sample slides:
    # slide 1 = title, 2 = divider, 3 = content, 4 = closing.
    template_slides = list(prs.slides)
    layout_title = template_slides[0].slide_layout if len(template_slides) >= 1 else find_layout(prs, ["title"])
    layout_section = template_slides[1].slide_layout if len(template_slides) >= 2 else find_layout(prs, ["divider", "section"])
    layout_content = template_slides[2].slide_layout if len(template_slides) >= 3 else find_layout(prs, ["content", "1 column", "title/subtitle"])
    layout_closing = template_slides[3].slide_layout if len(template_slides) >= 4 else find_layout(prs, ["closing", "thank"])
    layout_blank_light = find_layout(prs, ["light - blank", "blank", "title 1 column"])
    clear_existing_slides(prs)

    infra_eval_all = infra_assessment(model, ai_payload)
    plan = build_plan(model)
    for item in plan:
        t = item["type"]
        if t == "executive-summary":
            slide = prs.slides.add_slide(layout_content)
            s = model.summary()
            total_db = int(s.get("db_count") or 0)
            total_hosts = int(s.get("host_count") or 0)
            total_instances = int(s.get("instance_count") or 0)
            total_cohorts = len(model.cohorts())
            iops_items = model.iops_by_cohort()
            infra_eval = infra_eval_all

            set_title_and_subtitle(
                slide,
                "Executive summary: INTERNAL ORACLE - Not to share with customer",
                f"Scope: {total_db} DBs | {total_instances} instances | {total_hosts} hosts | {total_cohorts} cohorts",
            )
            max_iops = max((float(v) for _, v in iops_items), default=0.0)
            tier = str(infra_eval["tier"])
            conf = float(infra_eval["confidence"])
            req_vcpu = float(s.get("vcpu_total") or 0.0)
            req_ecpu = float(infra_eval.get("required_ecpu_total") or (req_vcpu * ECPU_PER_VCPU_EXADATA))
            req_storage_tb = float(s.get("allocated_storage_gb") or 0.0) / 1024.0
            req_memory_gb = float(s.get("memory_gb_total") or 0.0)
            req_iops = float(max_iops)

            add_exec_cards(
                slide,
                [
                    {"value": fmt_num(req_ecpu, 1), "label": "Required ECPUs (approx)", "sub": "from total vCPU", "color": ORACLE_COLORS["accent1"]},
                    {"value": f"{req_storage_tb:.1f} TB", "label": "Required Storage", "sub": "allocated footprint", "color": ORACLE_COLORS["accent2"]},
                    {"value": f"{req_memory_gb/1024.0:.1f} TB", "label": "Required Memory", "sub": f"{fmt_num(req_memory_gb,1)} GB", "color": ORACLE_COLORS["accent5"]},
                    {"value": fmt_num(req_iops, 0), "label": "Peak Cohort IOPS", "sub": "p95/peak indicator", "color": ORACLE_COLORS["accent3"]},
                    {"value": tier, "label": "Recommended Platform", "sub": f"conf {conf:.2f}", "color": ORACLE_COLORS["accent1"]},
                ],
            )

            add_section_label(slide, "Cohort Deployment Fit (Bottom-Up)", 0.79, 3.02)
            # Requirements summary matrix for deployment decisions.
            selected_rows = model.selected_rows()
            base_limits = infra_eval.get("base_limits") or {}
            xs_limits = infra_eval.get("xs_limits") or {}
            exa_limits = infra_eval.get("exa_env_limits") or {}
            exa_envs = max(1, int(infra_eval.get("exadata_envs") or 1))
            per_env_ecpu = req_ecpu / exa_envs
            per_env_mem = req_memory_gb / exa_envs
            per_env_storage = req_storage_tb / exa_envs
            per_env_iops = req_iops / exa_envs

            max_inst_cpu = max((float(r.init_cpu_count or r.logical_cpu_count or 0.0) for r in selected_rows), default=0.0)
            max_inst_mem = max((float(r.mem_gb or 0.0) for r in selected_rows), default=0.0)
            max_inst_sga = max((float(r.sga_gb or 0.0) for r in selected_rows), default=0.0)
            max_inst_pga = max((float(r.pga_gb or 0.0) for r in selected_rows), default=0.0)
            max_inst_storage = 0.0
            max_inst_iops = 0.0
            for r in selected_rows:
                st = model.db_metric_stats(r.cohort, r.db, "Allocated Storage (GB)") or {}
                max_inst_storage = max(max_inst_storage, float(st.get("p95") or st.get("max") or 0.0) / 1024.0)
                it = model.db_metric_stats(r.cohort, r.db, "DB IOPS") or {}
                max_inst_iops = max(max_inst_iops, float(it.get("p95") or it.get("max") or 0.0))

            max_cohort_vcpu = max((float(model.summary(model.rows_for_cohort(c)).get("vcpu_total") or 0.0) for c in model.cohorts()), default=0.0)
            max_cohort_ecpu = max_cohort_vcpu * ECPU_PER_VCPU_EXADATA
            max_cohort_storage = max((float(model.summary(model.rows_for_cohort(c)).get("allocated_storage_gb") or 0.0) / 1024.0 for c in model.cohorts()), default=0.0)

            cpu_series = model.cpu_series_global()
            all_cpu_points = [v for _, _, ys in cpu_series for v in ys]
            cpu_p95 = quantile(sorted(all_cpu_points), 0.95) if all_cpu_points else 0.0
            vcpu_usage_pct = pct(cpu_p95, req_vcpu)

            rac_flags = {model.db_is_rac(r.db) for r in selected_rows}
            rac_state = "Mixed" if len(rac_flags) > 1 else ("Yes" if True in rac_flags else "No")

            def fit_txt(req: float, lim: float, unit: str, dec: int = 1) -> str:
                st = "Fit" if req <= lim else "Not Fit"
                return f"{st} ({fmt_num(req, dec)}/{fmt_num(lim, dec)} {unit})"

            matrix_rows = [
                [
                    "VCPU total",
                    f"{fmt_num(req_vcpu,1)} vCPU",
                    f"Per-instance max: {fit_txt(max_inst_cpu, float(base_limits.get('ecpu', 0.0)), 'ECPU', 1)}",
                    f"Largest cohort: {fit_txt(max_cohort_ecpu, float(xs_limits.get('ecpu', 0.0)), 'ECPU', 1)}",
                    f"Per env: {fit_txt(per_env_ecpu, float(exa_limits.get('ecpu', 0.0)), 'ECPU', 1)}",
                ],
                [
                    "VCPU Usage %",
                    f"{vcpu_usage_pct:.1f}%",
                    "Advisory metric",
                    "Advisory metric",
                    "Advisory metric",
                ],
                [
                    "Memory",
                    f"{fmt_num(req_memory_gb,1)} GB",
                    f"Per-instance max: {fit_txt(max_inst_mem, float(base_limits.get('memory_gb', 0.0)), 'GB', 1)}",
                    "Resource-scaled",
                    f"Per env: {fit_txt(per_env_mem, float(exa_limits.get('memory_gb', 0.0)), 'GB', 1)}",
                ],
                [
                    "Storage",
                    f"{fmt_num(req_storage_tb,1)} TB",
                    f"Per-instance max: {fit_txt(max_inst_storage, float(base_limits.get('storage_tb', 0.0)), 'TB', 1)}",
                    f"Largest cohort: {fit_txt(max_cohort_storage, float(xs_limits.get('storage_tb', 0.0)), 'TB', 1)}",
                    f"Per env: {fit_txt(per_env_storage, float(exa_limits.get('storage_tb', 0.0)), 'TB', 1)}",
                ],
                [
                    "IOPS",
                    fmt_num(req_iops, 0),
                    f"Per-instance max: {fit_txt(max_inst_iops, float(base_limits.get('iops', 0.0)), 'IOPS', 0)}",
                    "Resource profile",
                    f"Per env: {fit_txt(per_env_iops, float(exa_limits.get('iops', 0.0)), 'IOPS', 0)}",
                ],
                [
                    "SGA",
                    f"{fmt_num(s.get('sga_total') or 0.0,1)} GB",
                    f"Max instance: {fmt_num(max_inst_sga,1)} GB",
                    "Included in memory profile",
                    "Included in env memory sizing",
                ],
                [
                    "PGA",
                    f"{fmt_num(s.get('pga_total') or 0.0,1)} GB",
                    f"Max instance: {fmt_num(max_inst_pga,1)} GB",
                    "Included in memory profile",
                    "Included in env memory sizing",
                ],
                [
                    "RAC",
                    rac_state,
                    "RAC instances count as 2 DB servers",
                    "One Exascale deployment per cohort",
                    "All cohorts consolidated in Exadata Dedicated",
                ],
            ]
            add_table(
                slide,
                headers=["Requirement", "Current Scope", "Base Database", "Exadata Exascale (ExaDB-XS)", "Exadata Dedicated (ExaDB-D)"],
                rows=matrix_rows,
                x=0.79,
                y=3.20,
                w=12.0,
                h=2.35,
                fit_cols={2, 3, 4},
            )
            add_panel(
                slide,
                "Warning",
                [
                    "INTERNAL ORACLE ONLY - do not share this slide with customer.",
                    "This recommendation is deterministic from AWR quantitative metrics.",
                    "Validate final positioning with Oracle solution architecture review.",
                ],
                x=0.79,
                y=5.62,
                w=12.0,
                h=1.0,
                dark=False,
                body_font_size=12,
            )
            if tier == "Base Database":
                sizing_line = f"Sizing target: {total_instances} Base Databases (one per instance)."
            elif tier == "Exascale":
                sizing_line = f"Sizing target: {total_cohorts} Exascale deployments (one per cohort), hosting {total_instances} instances."
            else:
                sizing_line = (
                    f"Sizing target: {infra_eval.get('exadata_envs', 1)} Exadata Dedicated environment(s), "
                    f"{infra_eval.get('exadata_db_servers_total', 0)} DB servers and "
                    f"{infra_eval.get('exadata_storage_servers_total', 0)} Storage servers."
                )
            add_footer_note(slide, sizing_line, x=0.8, y=6.90, w=11.9)
            continue

        if t == "deployment-base":
            slide = prs.slides.add_slide(layout_content)
            set_title_and_subtitle(slide, item["title"], item.get("subtitle", ""))
            base_limits = infra_eval_all.get("base_limits") or {}
            add_exec_cards(
                slide,
                [
                    {"value": "Per Instance", "label": "Sizing Scope", "sub": "Each DB instance is evaluated independently", "color": ORACLE_COLORS["accent1"]},
                    {"value": f"{fmt_num(base_limits.get('ecpu', 0),0)} ECPU", "label": "CPU Threshold", "sub": "Per instance", "color": ORACLE_COLORS["accent5"]},
                    {"value": f"{fmt_num(base_limits.get('memory_gb', 0),0)} GB", "label": "Memory Threshold", "sub": "Per instance", "color": ORACLE_COLORS["accent3"]},
                    {"value": f"{fmt_num(base_limits.get('storage_tb', 0),0)} TB", "label": "Storage Threshold", "sub": "Per instance", "color": ORACLE_COLORS["accent1"]},
                    {"value": f"{fmt_num(base_limits.get('iops', 0),0)} IOPS", "label": "IOPS Threshold", "sub": "Per instance", "color": ORACLE_COLORS["accent2"]},
                ],
            )
            add_section_label(slide, "Base Database Sizing Limits", 0.79, 3.02)
            map_rows = [
                ["CPU", f"{fmt_num(base_limits.get('ecpu',0),0)}", "ECPU", "Maximum CPU per instance"],
                ["Memory", f"{fmt_num(base_limits.get('memory_gb',0),0)}", "GB", "Maximum memory per instance"],
                ["Allocated Storage", f"{fmt_num(base_limits.get('storage_tb',0),0)}", "TB", "Maximum allocated storage per instance"],
                ["IOPS", f"{fmt_num(base_limits.get('iops',0),0)}", "IOPS", "Maximum IOPS per instance"],
                ["RAC Rule", "2 DB servers", "count", "If an instance is RAC, count two DB servers in inventory"],
            ]
            add_table(
                slide,
                headers=["Parameter", "Limit", "Unit", "Notes"],
                rows=map_rows,
                x=0.79,
                y=3.20,
                w=12.0,
                h=3.55,
            )
            add_footer_note(
                slide,
                "[A] Citations: oracle-base-database-service.pdf | p.30 | Table 1-2 Flexible Shapes; p.33 | Table 1-4 Available data storage; p.7 | Service limits.",
                x=0.8,
                y=6.93,
                w=11.9,
            )
            continue

        if t == "deployment-exascale":
            slide = prs.slides.add_slide(layout_content)
            set_title_and_subtitle(slide, item["title"], item.get("subtitle", ""))
            xs_limits = infra_eval_all.get("xs_limits") or {}
            add_exec_cards(
                slide,
                [
                    {"value": "Per Cohort", "label": "Sizing Scope", "sub": "One Exascale deployment per cohort", "color": ORACLE_COLORS["accent1"]},
                    {"value": "vCPU x 2", "label": "CPU Conversion", "sub": "Exadata-family conversion", "color": ORACLE_COLORS["accent2"]},
                    {"value": "8", "label": "Min ECPU / VM", "sub": "Capacity limit", "color": ORACLE_COLORS["accent5"]},
                    {"value": "200", "label": "Max ECPU / VM", "sub": "Capacity limit", "color": ORACLE_COLORS["accent3"]},
                    {"value": "10", "label": "Max VMs / Cluster", "sub": "Capacity limit", "color": ORACLE_COLORS["accent1"]},
                ],
            )
            add_section_label(slide, "Exascale Sizing Limits", 0.79, 3.02)
            xs_rows = [
                ["Minimum VM cluster size", "Single-node VM cluster", "cluster", "Baseline deployment minimum"],
                ["Minimum ECPUs per VM", "8", "ECPU/VM", "Capacity lower bound"],
                ["Maximum ECPUs per VM", "200", "ECPU/VM", "Capacity upper bound"],
                ["Maximum VMs in VM cluster", "10", "VMs", "Cluster size upper bound"],
                ["File system storage per VM (min)", "220 (26ai) / 260 (19c)", "GB/VM", "Minimum billed capacity"],
                ["File system storage per VM (max)", "2", "TB/VM", "Capacity upper bound"],
                ["Exascale Vault DB storage (min)", "300", "GB/VM cluster", "Minimum cluster database storage"],
                ["Memory scaling granularity", "2.75", "GB per total ECPU", "Scaling increment"],
                ["Current script cohort thresholds", f"{fmt_num(xs_limits.get('ecpu',0),0)} ECPU and {fmt_num(xs_limits.get('storage_tb',0),0)} TB", "per cohort", "Internal policy threshold"],
            ]
            add_table(
                slide,
                headers=["Parameter", "Limit", "Unit", "Notes"],
                rows=xs_rows,
                x=0.79,
                y=3.20,
                w=12.0,
                h=3.55,
            )
            add_footer_note(
                slide,
                "[B] Citations: Getting Started with Oracle Exadata Database Service on Exascale Infrastructure Deployment.pdf | p.30-p31 | Capacity limits and VM storage capacities. oracle-exadata-exascale-users-guide-exscl.pdf | p.31-p32 | Resource management context.",
                x=0.8,
                y=6.93,
                w=11.9,
            )
            continue

        if t == "deployment-exadata":
            slide = prs.slides.add_slide(layout_content)
            set_title_and_subtitle(slide, item["title"], item.get("subtitle", ""))
            exa_limits = infra_eval_all.get("exa_env_limits") or {}
            add_exec_cards(
                slide,
                [
                    {"value": "vCPU x 2", "label": "CPU Conversion", "sub": "Exadata-family ECPU conversion", "color": ORACLE_COLORS["accent1"]},
                    {"value": f"{fmt_num(exa_limits.get('ecpu', 0),0)}", "label": "Usable ECPUs / System", "sub": "X11M minimum config reference", "color": ORACLE_COLORS["accent2"]},
                    {"value": f"{fmt_num(exa_limits.get('memory_gb', 0),0)} GB", "label": "Memory Limit / Cohort", "sub": "70% of one DB server memory", "color": ORACLE_COLORS["accent5"]},
                    {"value": f"{fmt_num(exa_limits.get('storage_tb', 0),0)} TB", "label": "Storage Envelope", "sub": "Sizing policy upper envelope", "color": ORACLE_COLORS["accent3"]},
                    {"value": f"{fmt_num(exa_limits.get('iops', 0),0)}", "label": "IOPS Envelope", "sub": "Sizing policy upper envelope", "color": ORACLE_COLORS["accent1"]},
                ],
            )
            add_section_label(slide, "Exadata Dedicated Sizing Limits", 0.79, 3.02)
            exa_rows = [
                [
                    "Total usable ECPUs in DB servers per system",
                    f"{fmt_num(exa_limits.get('ecpu',0),0)}",
                    "ECPU/system",
                    "Minimum configuration property",
                ],
                [
                    "Cohort memory limit (policy)",
                    f"{fmt_num(exa_limits.get('memory_gb',0),0)}",
                    "GB/cohort",
                    "70% of one DB server memory (1390 GB)",
                ],
                [
                    "Starting infrastructure footprint",
                    "2 DB servers + 3 Storage servers",
                    "count",
                    "Baseline in X11M dedicated model",
                ],
                [
                    "Database server scale range",
                    "2 to 32",
                    "DB servers",
                    "Horizontal scale options",
                ],
                [
                    "Storage server scale range",
                    "3 to 64",
                    "Storage servers",
                    "Horizontal scale options",
                ],
                [
                    "ECPU activation increment",
                    "In multiples based on DB server count",
                    "ECPU step",
                    "Operational scaling granularity",
                ],
                [
                    "Minimum enabled ECPUs",
                    "0",
                    "ECPU",
                    "Minimum (default) configuration property",
                ],
            ]
            add_table(
                slide,
                headers=["Parameter", "Limit", "Unit", "Notes"],
                rows=exa_rows,
                x=0.79,
                y=3.20,
                w=12.0,
                h=3.55,
            )
            add_footer_note(
                slide,
                "[C] Citations: exadata-database-service-dedicated-infrastructure-administrators-guide.pdf | p.31 | Scaling operations; p.33 | Exadata shape configuration; p.34 | ECPU and memory limits. Policy applied: cohort memory cap = 70% of one DB server memory (1390 GB => 973 GB).",
                x=0.8,
                y=6.93,
                w=11.9,
            )
            continue

        if t == "title":
            slide = prs.slides.add_slide(layout_title)
            customer = report_meta.get("customerName") or ""
            subtitle = f"AWR Analysis for: {customer}" if customer else "AWR Analysis"
            set_title_and_subtitle(slide, item["title"], subtitle)

            # Fill template placeholders:
            # "Name" -> "Prepared by:"
            # "Presenter's Title" -> Sales Rep / Architect / Engineer names (bold names).
            # Some templates merge both labels into one textbox, so detect both separately
            # and also pick a role box by position when labels are not directly matchable.
            name_placeholder = None
            presenter_placeholder = None
            combined_placeholder = None
            role_box_fallback = None
            best_dist = None
            target_left = Cm(2.2)
            target_top = Cm(13.56)
            for shp in slide.shapes:
                if not getattr(shp, "has_text_frame", False):
                    continue
                if shp == slide.shapes.title:
                    continue
                txt = (shp.text_frame.text or "").strip()
                norm = re.sub(r"[^a-z]+", " ", txt.lower()).strip()
                has_name = bool(re.search(r"\bname\b", norm))
                has_presenter = bool(re.search(r"\bpresenter\b", norm))
                has_title = bool(re.search(r"\btitle\b", norm))
                if has_name and has_presenter:
                    combined_placeholder = shp
                if has_name and name_placeholder is None:
                    name_placeholder = shp
                if has_presenter and (has_title or "presenter s" in norm) and presenter_placeholder is None:
                    presenter_placeholder = shp
                # Position fallback for role box.
                try:
                    dx = abs(int(shp.left) - int(target_left))
                    dy = abs(int(shp.top) - int(target_top))
                    dist = dx + dy
                    if best_dist is None or dist < best_dist:
                        best_dist = dist
                        role_box_fallback = shp
                except Exception:
                    pass

            role_items = [
                ("Sales Rep", report_meta.get("salesRepName") or ""),
                ("Architect", report_meta.get("architectName") or ""),
                ("Engineer", report_meta.get("engineerName") or ""),
            ]
            role_items = [(r, n) for (r, n) in role_items if n]

            # Prefer true combined placeholder when available.
            if combined_placeholder is not None:
                role_box = combined_placeholder
            elif presenter_placeholder is not None:
                role_box = presenter_placeholder
            else:
                role_box = role_box_fallback

            # If combined placeholder is used, write "Prepared by:" + role lines in one box.
            if role_box is not None and role_box == combined_placeholder:
                tf = role_box.text_frame
                tf.clear()
                p0 = tf.paragraphs[0]
                p0.text = "Prepared by:"
                p0.font.name = "Oracle Sans Tab"
                p0.font.size = Pt(16)
                p0.font.color.rgb = ORACLE_COLORS["text"]
                p0.alignment = PP_ALIGN.LEFT
                for i, (role, name) in enumerate(role_items):
                    p = tf.add_paragraph()
                    p.text = ""
                    run_role = p.add_run()
                    run_role.text = f"{role}: "
                    run_role.font.name = "Oracle Sans Tab"
                    run_role.font.size = Pt(16)
                    run_role.font.bold = False
                    run_role.font.color.rgb = ORACLE_COLORS["text"]
                    run_name = p.add_run()
                    run_name.text = name
                    run_name.font.name = "Oracle Sans Tab"
                    run_name.font.size = Pt(16)
                    run_name.font.bold = True
                    run_name.font.color.rgb = ORACLE_COLORS["text"]
                    p.alignment = PP_ALIGN.LEFT
            else:
                if name_placeholder is not None:
                    ntf = name_placeholder.text_frame
                    ntf.clear()
                    p = ntf.paragraphs[0]
                    p.text = "Prepared by:"
                    p.font.name = "Oracle Sans Tab"
                    p.font.size = Pt(16)
                    p.font.color.rgb = ORACLE_COLORS["text"]
                    p.alignment = PP_ALIGN.LEFT
                if presenter_placeholder is not None:
                    ptf = presenter_placeholder.text_frame
                    ptf.clear()
                    for i, (role, name) in enumerate(role_items):
                        p = ptf.paragraphs[0] if i == 0 else ptf.add_paragraph()
                        p.text = ""
                        run_role = p.add_run()
                        run_role.text = f"{role}: "
                        run_role.font.name = "Oracle Sans Tab"
                        run_role.font.size = Pt(16)
                        run_role.font.bold = False
                        run_role.font.color.rgb = ORACLE_COLORS["text"]
                        run_name = p.add_run()
                        run_name.text = name
                        run_name.font.name = "Oracle Sans Tab"
                        run_name.font.size = Pt(16)
                        run_name.font.bold = True
                        run_name.font.color.rgb = ORACLE_COLORS["text"]
                        p.alignment = PP_ALIGN.LEFT
                if name_placeholder is None and presenter_placeholder is None:
                    # Last-resort fallback if template role box is not detectable.
                    rb = slide.shapes.add_textbox(Cm(2.2), Cm(13.56), Cm(15.8), Cm(3.2))
                    tf = rb.text_frame
                    tf.clear()
                    p0 = tf.paragraphs[0]
                    p0.text = "Prepared by:"
                    p0.font.name = "Oracle Sans Tab"
                    p0.font.size = Pt(16)
                    p0.font.color.rgb = ORACLE_COLORS["text"]
                    p0.alignment = PP_ALIGN.LEFT
                    for role, name in role_items:
                        p = tf.add_paragraph()
                        p.text = ""
                        run_role = p.add_run()
                        run_role.text = f"{role}: "
                        run_role.font.name = "Oracle Sans Tab"
                        run_role.font.size = Pt(16)
                        run_role.font.bold = False
                        run_role.font.color.rgb = ORACLE_COLORS["text"]
                        run_name = p.add_run()
                        run_name.text = name
                        run_name.font.name = "Oracle Sans Tab"
                        run_name.font.size = Pt(16)
                        run_name.font.bold = True
                        run_name.font.color.rgb = ORACLE_COLORS["text"]
                        p.alignment = PP_ALIGN.LEFT

            gen = slide.shapes.add_textbox(Cm(2.2), Cm(18.0), Cm(10.5), Cm(0.8))
            gtf = gen.text_frame
            gtf.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
            gtf.paragraphs[0].font.name = "Oracle Sans Tab"
            gtf.paragraphs[0].font.size = Pt(12)
            gtf.paragraphs[0].font.color.rgb = ORACLE_COLORS["muted"]
            gtf.paragraphs[0].alignment = PP_ALIGN.LEFT
            continue

        if t == "separator":
            slide = prs.slides.add_slide(layout_section)
            set_title_and_subtitle(slide, item["title"], item.get("subtitle", ""))
            continue

        if t == "summary":
            slide = prs.slides.add_slide(layout_content)
            s = model.summary()
            versions = [(k, v) for k, v in model.version_counts()]
            top_ver = versions[0] if versions else ("N/A", 0)
            total_ver = sum(v for _, v in versions) or 1
            cpu_series = model.cpu_series_global()
            all_cpu_points = [v for _, _, ys in cpu_series for v in ys]
            peak_cpu = max(all_cpu_points) if all_cpu_points else None
            p95_cpu = None
            if all_cpu_points:
                sorted_cpu = sorted(all_cpu_points)
                p95_cpu = quantile(sorted_cpu, 0.95)
            total_db = int(s["db_count"])
            total_hosts = int(s["host_count"])
            total_instances = int(s["instance_count"])
            total_cohorts = len(model.cohorts())
            freshness = model.freshness_ranges()
            cpu_cohort_stats = model.cpu_stats_by_cohort()
            mem_items = model.memory_by_cohort()
            iops_items = model.iops_by_cohort()

            title_msg = "Summary of AWR Capture."
            subtitle_msg = f"Instances: {total_instances} (Databases: {total_db}, Hosts: {total_hosts})"
            set_title_and_subtitle(slide, title_msg, subtitle_msg)

            cpu_value = f"{peak_cpu:.1f} vCPU" if peak_cpu is not None else "N/A"
            cpu_sub = f"P95: {p95_cpu:.1f} vCPU" if p95_cpu is not None else "P95: N/A"
            add_exec_cards(
                slide,
                [
                    {
                        "value": fmt_num(total_instances, 0),
                        "label": "Total Instances",
                        "sub": f"{fmt_num(total_db,0)} databases",
                        "color": ORACLE_COLORS["accent1"],
                    },
                    {
                        "value": fmt_num(total_hosts, 0),
                        "label": "Total Hosts",
                        "sub": f"{fmt_num(total_cohorts,0)} cohorts",
                        "color": ORACLE_COLORS["accent1"],
                    },
                    {
                        "value": f"{s['memory_gb_total']/1024.0:.1f} TB",
                        "label": "Total Memory",
                        "sub": f"{fmt_num(s['memory_gb_total'],1)} GB",
                        "color": ORACLE_COLORS["accent5"],
                    },
                    {
                        "value": fmt_num(s["vcpu_total"], 1),
                        "label": "Total vCPU",
                        "sub": "Configured capacity",
                        "color": ORACLE_COLORS["accent3"],
                    },
                    {
                        "value": f"{s['allocated_storage_gb']/1024.0:.1f} TB",
                        "label": "Allocated Storage",
                        "sub": f"{fmt_num(s['allocated_storage_gb'],1)} GB",
                        "color": ORACLE_COLORS["accent2"],
                    },
                ],
            )
            versions_lines = []
            top_versions = versions[:4]
            for ver, count in top_versions:
                pct_val = 100.0 * count / max(total_db, 1)
                versions_lines.append(f"{ver}: {count}/{total_db} DBs ({pct_val:.1f}%)")
            if len(versions) > 4:
                others = sum(c for _, c in versions[4:])
                pct_val = 100.0 * others / max(total_db, 1)
                versions_lines.append(f"Others: {others}/{total_db} DBs ({pct_val:.1f}%)")
            add_big_versions_card(
                slide,
                "Database Versions",
                versions_lines or ["No version data available"],
                x=8.6,
                y=3.15,
                w=4.2,
                h=1.58,
            )

            learned = summary_dba_learnings(
                model=model,
                s=s,
                total_db=total_db,
                total_instances=total_instances,
                total_hosts=total_hosts,
                total_cohorts=total_cohorts,
                versions=versions,
                cpu_stats=cpu_cohort_stats,
                iops_items=iops_items,
            )
            add_panel(slide, "What We Learned", learned, x=0.79, y=3.15, w=7.6, h=1.58, dark=False, body_font_size=14)
            c_y = 5.05
            c_h = 1.65
            c_gap = 0.12
            c_w = (12.0 - (3 * c_gap)) / 4.0
            c1_x = 0.79
            c2_x = c1_x + c_w + c_gap
            c3_x = c2_x + c_w + c_gap
            c4_x = c3_x + c_w + c_gap

            add_section_label(slide, "AWR capture timeline", c1_x, 4.88)
            add_section_label(slide, "Memory", c2_x, 4.88)
            add_section_label(slide, "IOPS", c3_x, 4.88)
            add_section_label(slide, "CPU Box Plot", c4_x, 4.88)
            if freshness:
                add_freshness_timeline_chart(slide, freshness, x=c1_x, y=c_y, w=c_w, h=c_h, legend_font_size=4)
            else:
                add_freshness_unavailable_note(slide, x=c1_x, y=5.62, w=c_w)
            add_bar_chart(slide, mem_items, x=c2_x, y=c_y, w=c_w, h=c_h, series_name="Memory (GB)", x_tick_font_size=4)
            add_bar_chart(slide, iops_items, x=c3_x, y=c_y, w=c_w, h=c_h, series_name="DB IOPS", x_tick_font_size=4)
            add_cohort_cpu_boxplot_chart(slide, cpu_cohort_stats, x=c4_x, y=c_y, w=c_w, h=c_h)
            add_footer_note(
                slide,
                "Source: instances, database_statistics, plots_html | Metrics shown for selected PPT scope",
                x=0.8,
                y=6.92,
                w=11.9,
            )
            continue

        if t == "cohort":
            c = item["cohort"]
            slide = prs.slides.add_slide(layout_content)
            rs = model.rows_for_cohort(c)
            sm = model.summary(rs)
            set_title_and_subtitle(slide, item["title"], f"Databases: {int(sm['db_count'])} | Instances: {int(sm['instance_count'])}")
            version_dbs: dict[str, set[str]] = defaultdict(set)
            for r in rs:
                version_dbs[model.db_versions.get(r.db, "Unknown")].add(r.db)
            versions = sorted(((ver, len(dbs)) for ver, dbs in version_dbs.items()), key=lambda kv: (-kv[1], kv[0]))
            top_ver = versions[0] if versions else ("Unknown", 0)
            total_db = int(sm["db_count"])
            total_hosts = int(sm["host_count"])
            total_instances = int(sm["instance_count"])

            add_exec_cards(
                slide,
                [
                    {
                        "value": f"{total_instances:d}",
                        "label": "Total Instances",
                        "sub": f"Across {total_db:d} DBs",
                        "color": ORACLE_COLORS["accent1"],
                    },
                    {
                        "value": f"{total_hosts:d}",
                        "label": "Total Hosts",
                        "sub": f"{total_instances:d} instances distributed",
                        "color": ORACLE_COLORS["accent3"],
                    },
                    {
                        "value": f"{sm['memory_gb_total']/1024.0:.1f} TB",
                        "label": "Total Memory",
                        "sub": f"{fmt_num(sm['memory_gb_total'],1)} GB",
                        "color": ORACLE_COLORS["accent2"],
                    },
                    {
                        "value": f"{fmt_num(sm['vcpu_total'],1)}",
                        "label": "Total vCPU",
                        "sub": "Configured capacity",
                        "color": ORACLE_COLORS["accent1"],
                    },
                    {
                        "value": f"{sm['allocated_storage_gb']/1024.0:.1f} TB",
                        "label": "Allocated Storage",
                        "sub": f"{fmt_num(sm['allocated_storage_gb'],1)} GB",
                        "color": ORACLE_COLORS["accent2"],
                    },
                ],
            )
            versions_lines = []
            top_versions = versions[:4]
            for ver, count in top_versions:
                pct_val = 100.0 * count / max(total_db, 1)
                versions_lines.append(f"{ver}: {count}/{total_db} DBs ({pct_val:.1f}%)")
            if len(versions) > 4:
                others = sum(cn for _, cn in versions[4:])
                pct_val = 100.0 * others / max(total_db, 1)
                versions_lines.append(f"Others: {others}/{total_db} DBs ({pct_val:.1f}%)")
            add_big_versions_card(
                slide,
                "Database Versions",
                versions_lines or ["No version data available"],
                x=8.6,
                y=3.15,
                w=4.2,
                h=1.58,
            )

            cohort_freshness = model.freshness_ranges_by_instance_in_cohort(c)
            mem_items_c = model.metric_by_db(c, "DB Memory (MB)", scale=(1.0 / 1024.0))
            iops_items_c = model.metric_by_db(c, "DB IOPS", scale=1.0)
            inst_stats = []
            for inst_name, _xs, ys in model.cpu_series_instances_in_cohort(c):
                st = stats_from_values(ys)
                if st:
                    inst_stats.append((inst_name, st))

            learned = cohort_dba_learnings(
                model=model,
                cohort=c,
                rs=rs,
                sm=sm,
                versions=versions,
                inst_stats=inst_stats,
                mem_items_c=mem_items_c,
                iops_items_c=iops_items_c,
            )
            add_panel(slide, "What We Learned", learned, x=0.79, y=3.15, w=7.6, h=1.58, dark=False, body_font_size=14)

            c_y = 5.05
            c_h = 1.65
            c_gap = 0.12
            c_w = (12.0 - (3 * c_gap)) / 4.0
            c1_x = 0.79
            c2_x = c1_x + c_w + c_gap
            c3_x = c2_x + c_w + c_gap
            c4_x = c3_x + c_w + c_gap

            add_section_label(slide, "AWR capture timeline", c1_x, 4.88)
            add_section_label(slide, "Memory", c2_x, 4.88)
            add_section_label(slide, "IOPS", c3_x, 4.88)
            add_section_label(slide, "CPU Box Plot", c4_x, 4.88)
            if cohort_freshness:
                add_freshness_timeline_chart(slide, cohort_freshness, x=c1_x, y=c_y, w=c_w, h=c_h, legend_font_size=4)
            else:
                add_freshness_unavailable_note(slide, x=c1_x, y=5.62, w=c_w)
            add_bar_chart(slide, mem_items_c, x=c2_x, y=c_y, w=c_w, h=c_h, series_name="Memory (GB)", x_tick_font_size=4)
            add_bar_chart(slide, iops_items_c, x=c3_x, y=c_y, w=c_w, h=c_h, series_name="DB IOPS", x_tick_font_size=4)
            add_cohort_cpu_boxplot_chart(slide, inst_stats, x=c4_x, y=c_y, w=c_w, h=c_h)
            continue

        if t == "instance":
            slide = prs.slides.add_slide(layout_content)
            c = item["cohort"]
            inst = item["instance"]
            db = item["db"]
            rs = [r for r in model.rows_for_cohort(c) if r.instance == inst]
            row = rs[0] if rs else None
            db_name = model.display_db_name(db)
            db_type, pdb_count = model.db_type_and_pdb_count(db)
            db_version = model.db_version(db)
            is_rac = "Yes" if model.db_is_rac(db) else "No"
            set_title_and_subtitle(slide, item["title"], item["subtitle"])
            add_kpi_boxes(
                slide,
                [
                    ("Cohort", c),
                    ("Database", db_name),
                    ("DB Version", db_version),
                    ("DB Type", db_type),
                    ("PDBs", pdb_count),
                    ("RAC", is_rac),
                    ("Host", row.host if row else "N/A"),
                    ("vCPU", fmt_num((row.init_cpu_count if row else 0), 1)),
                    ("Memory (GB)", fmt_num((row.mem_gb if row else 0), 1)),
                    ("SGA (GB)", fmt_num((row.sga_gb if row else 0), 1)),
                    ("PGA (GB)", fmt_num((row.pga_gb if row else 0), 1)),
                ],
            )

            mem_title = slide.shapes.add_textbox(Inches(0.79), Inches(3.40), Inches(6.2), Inches(0.24))
            mem_tf = mem_title.text_frame
            mem_tf.clear()
            mem_p = mem_tf.paragraphs[0]
            mem_p.text = "Memory Consumption"
            mem_p.font.name = "Oracle Sans Tab"
            mem_p.font.bold = True
            mem_p.font.size = Pt(14)
            mem_p.font.color.rgb = ORACLE_COLORS["muted"]
            mem_p.alignment = PP_ALIGN.LEFT

            host_mem = to_num(row.mem_gb if row else 0.0)
            inst_mem = to_num((row.sga_gb if row else 0.0) + (row.pga_gb if row else 0.0))
            cpu_inst_series = [x for x in model.cpu_series_instances_in_cohort(c) if x[0] == inst]
            xs_mem = cpu_inst_series[0][1] if cpu_inst_series else (model.cpu_by_cohort.get(c)[0] if model.cpu_by_cohort.get(c) else [])
            mem_series = []
            if xs_mem:
                mem_series = [
                    ("Server Memory (GB)", xs_mem, [host_mem] * len(xs_mem)),
                    ("Instance Memory (GB)", xs_mem, [inst_mem] * len(xs_mem)),
                ]
            add_line_chart(
                slide,
                mem_series,
                x=0.79,
                y=3.58,
                w=8.05,
                h=1.65,
                show_legend=True,
                y_min=None,
                y_max=None,
                x_tick_font_size=4,
                y_tick_font_size=8,
            )
            mem_stats = model.db_metric_stats(c, db, "DB Memory (MB)")
            add_vcpu_boxplot_chart(
                slide,
                mem_stats,
                x=8.95,
                y=3.40,
                w=3.84,
                h=1.83,
                title_text="Memory Consumption (Box Plot)",
                value_scale=(1.0 / 1024.0),
            )

            cpu_title = slide.shapes.add_textbox(Inches(0.79), Inches(5.42), Inches(4.6), Inches(0.24))
            cpu_tf = cpu_title.text_frame
            cpu_tf.clear()
            cpu_p = cpu_tf.paragraphs[0]
            cpu_p.text = "CPU Consumption"
            cpu_p.font.name = "Oracle Sans Tab"
            cpu_p.font.bold = True
            cpu_p.font.size = Pt(14)
            cpu_p.font.color.rgb = ORACLE_COLORS["muted"]
            cpu_p.alignment = PP_ALIGN.LEFT
            add_line_chart(
                slide,
                cpu_inst_series,
                x=0.79,
                y=5.60,
                w=8.05,
                h=1.65,
                show_legend=False,
                y_min=None,
                y_max=None,
                x_tick_font_size=4,
                y_tick_font_size=8,
                chart_kind="area",
                line_width_pt=2.0,
                area_transparency=0.9,
            )
            # Use the same sampled instance series as the line chart to keep boxplot and line fully consistent.
            vcpu_stats = None
            if cpu_inst_series and cpu_inst_series[0][1] and cpu_inst_series[0][2]:
                _sx, _sy = downsample(cpu_inst_series[0][1], cpu_inst_series[0][2], 36)
                vcpu_stats = stats_from_values(_sy)
            add_vcpu_boxplot_chart(
                slide,
                vcpu_stats,
                x=8.95,
                y=5.42,
                w=3.84,
                h=1.83,
                title_text="CPU Consumption (Box Plot)",
                value_scale=1.0,
            )
            add_footer_note(slide, instance_dba_comment(model, c, db, inst, row), x=0.8, y=1.50, w=11.9)
            continue

        if t == "closing":
            slide = prs.slides.add_slide(layout_closing)
            set_title_and_subtitle(slide, item["title"], item.get("subtitle", ""))
            continue

        # Safety fallback
        slide = prs.slides.add_slide(layout_blank_light)
        set_title_and_subtitle(slide, item.get("title", "Slide"), item.get("subtitle", ""))

    prs.save(str(output_path))


def main():
    ap = argparse.ArgumentParser(description="Export AWR analysis PPT using Oracle template white layouts.")
    ap.add_argument("--input", required=True, help="Path to raw AWR JSON or saved report JSON.")
    ap.add_argument("--template", required=True, help="Path to Oracle .potx template.")
    ap.add_argument("--output", required=True, help="Output .pptx path.")
    args = ap.parse_args()
    export_ppt(Path(args.input), Path(args.template), Path(args.output))
    print(args.output)


if __name__ == "__main__":
    main()
